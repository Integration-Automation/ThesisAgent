"""Second clarity pass on exports/fang2026disentangling-zh-tw.pptx.

User feedback after pass 1: "解釋的不夠清楚" — pass 1 only glossed the
most obvious acronyms (ADA / VAE / DPI / Fiedler / EC / AT / za / zb)
and missed several first-use definitions plus the table-cell prose ``;``
sweep (the earlier semicolon audit walked only text frames, not table
cells inside `<a:graphicFrame>` shapes).

This pass adds:

  A) First-use definitions that pass 1 missed:
     - APD itself — never defined in full anywhere in the deck
     - SOTA — appears on slide 4 with no gloss
     - 越獄 / jailbreak — used as benchmark-category name without
       explaining what kind of attack it is
     - pp — appears in KPI as "-9.6 pp" with no expansion
     - ablation — appears in KPI ablation pair line
     - Transformer — first use slide 4 body 3
     - 知識蒸餾 — first standalone use slide 4 body 3
     - 譜分析 — first use slide 4 body 2 in a parenthetical
     - 改述 / 混淆 attack types — slide 3 pain-point bullet

  B) Slide 5 KPI labels — append English abbreviations in parens so
     subsequent slides' ``ADA / HOR / FPR / IL`` reference back cleanly:
       對抗偵測準確率 → 對抗偵測準確率 (ADA)
       有害輸出減量    → 有害輸出減量 (HOR)
       偽陽性率        → 偽陽性率 (FPR)
       推論延遲        → 推論延遲 (IL)

  C) Slide 6 technique-table cells — 5 prose ``;`` → ``,`` matching the
     "use ， not ;" convention added to CLAUDE.md this session; plus a
     brief gloss for RoBERTa (used unexplained).

Both text-frame shapes and table cells are walked. Idempotent: each
edit's replacement is pre-checked.
"""
from __future__ import annotations

import sys
from copy import deepcopy
from pathlib import Path

_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

SRC = Path("exports/fang2026disentangling-zh-tw.pptx")


# (slide_1based, find_substring, replace_substring, short_label)
EDITS: list[tuple[int, str, str, str]] = [
    # ===== Slide 3 — attack-type glosses =====================================
    (
        3,
        "規則式過濾遇到改述 / 混淆攻擊就破功",
        "規則式過濾遇到改述(paraphrasing,攻擊者改寫提示文字) / 混淆"
        "(obfuscation,以異常編碼或拼字繞過)攻擊就破功",
        "S3: gloss 改述 / 混淆 attack types",
    ),

    # ===== Slide 4 — subhead 4: gloss 越獄 at first use =======================
    (
        4,
        "4. 在三組越獄 benchmark 上的全面評估",
        "4. 在三組越獄(jailbreak,誘騙 LLM 繞過安全限制之提示)"
        "benchmark 上的全面評估",
        "S4 subhead4: gloss 越獄/jailbreak",
    ),

    # ===== Slide 4 — body 2: gloss 譜分析 =====================================
    (
        4,
        "在 za 的語意鄰域建圖,以譜分析(Fiedler 向量,圖 Laplacian 的第二小特徵向量,以及更高階特徵值)",
        "在 za 的語意鄰域建圖,以譜分析(spectral analysis,從圖的特徵值結構抓出社群分離模式)"
        "之 Fiedler 向量(圖 Laplacian 的第二小特徵向量)與更高階特徵值",
        "S4 body2: gloss 譜分析 + restructure Fiedler clause",
    ),

    # ===== Slide 4 — body 3: gloss Transformer + 知識蒸餾 =====================
    (
        4,
        "Transformer-based 對抗意圖偵測器,經知識蒸餾後每筆 12.3 ms",
        "Transformer-based(以注意力機制為核心的深度學習網路)對抗意圖偵測器,"
        "經知識蒸餾(把大型 teacher 模型的推理能力轉到小型 student 模型)"
        "後每筆 12.3 ms",
        "S4 body3: gloss Transformer + 知識蒸餾",
    ),

    # ===== Slide 4 — body 4: gloss APD + SOTA =================================
    (
        4,
        "JailBreakBench、ToxicPrompts、AdvPromptGen — APD 平均 ADA 92.3%,比 四個 SOTA 防禦至少高 5.6 個百分點。",
        "JailBreakBench、ToxicPrompts、AdvPromptGen — APD"
        "(Adversarial Prompt Disentanglement,對抗提示語意分解)平均 ADA 92.3%,"
        "比四個 SOTA(State-of-the-Art,當前最佳)防禦至少高 5.6 個百分點。",
        "S4 body4: gloss APD (full) + SOTA",
    ),

    # ===== Slide 5 — KPI labels: append English abbreviations =================
    # The KPI block lives in a single shape with multiple paragraphs.
    (
        5,
        "對抗偵測準確率: 92.3%",
        "對抗偵測準確率 (ADA): 92.3%",
        "S5: label → 對抗偵測準確率 (ADA)",
    ),
    (
        5,
        "有害輸出減量: 87.4%",
        "有害輸出減量 (HOR): 87.4%",
        "S5: label → 有害輸出減量 (HOR)",
    ),
    (
        5,
        "偽陽性率: 3.7%",
        "偽陽性率 (FPR): 3.7%",
        "S5: label → 偽陽性率 (FPR)",
    ),
    (
        5,
        "推論延遲: 12.3 ms",
        "推論延遲 (IL): 12.3 ms",
        "S5: label → 推論延遲 (IL)",
    ),
    (
        5,
        "移除 VAE 後 ADA 下降: -9.6 pp",
        "移除 VAE 後 ADA 下降: -9.6 pp(percentage points,百分點)",
        "S5: gloss pp = percentage points",
    ),
    (
        5,
        "(對照組: ablation:92.3 → 82.7)",
        "(對照組: ablation 消融實驗:92.3 → 82.7)",
        "S5: gloss ablation = 消融實驗",
    ),

    # ===== Slide 6 — technique-table cells: ; → , + RoBERTa gloss =============
    # Table cells are walked separately from text-frame shapes.
    (
        6,
        "關鍵字 pattern;65.4% ADA / 10.8 m",
        "關鍵字 pattern,65.4% ADA / 10.8 ms",
        "S6 cell: ; → , (also fix '10.8 m' truncation → '10.8 ms')",
    ),
    (
        6,
        "RoBERTa 掃輸出;84.1% / 45.6 ms",
        "RoBERTa(BERT 改良版預訓練語言模型)掃輸出,84.1% / 45.6 ms",
        "S6 cell: gloss RoBERTa + ; → ,",
    ),
    (
        6,
        "LLM fine-tune adv+benign;86.7%",
        "LLM fine-tune adv+benign,86.7%",
        "S6 cell: ; → ,",
    ),
    (
        6,
        "embedding 異常偵測;78.6% / 15.6 ms",
        "embedding 異常偵測,78.6% / 15.6 ms",
        "S6 cell: ; → ,",
    ),
    (
        6,
        "VAE + 譜圖 + 蒸餾 AID;92.3% / 12.3",
        "VAE + 譜圖 + 蒸餾 AID,92.3% / 12.3 ms",
        "S6 cell: ; → , (also fix '12.3' truncation → '12.3 ms')",
    ),
]


def _rewrite_paragraph_as_single_run(para, new_text: str) -> None:
    """Replace paragraph runs with a single run carrying ``new_text``,
    preserving rPr (font + colour) from the first original run."""
    p_elem = para._p
    first_r = p_elem.find(qn("a:r"))
    rpr_clone = None
    if first_r is not None:
        rpr_elem = first_r.find(qn("a:rPr"))
        if rpr_elem is not None:
            rpr_clone = deepcopy(rpr_elem)
    for r in p_elem.findall(qn("a:r")):
        p_elem.remove(r)
    new_r = p_elem.makeelement(qn("a:r"), {})
    if rpr_clone is not None:
        new_r.append(rpr_clone)
    new_t = p_elem.makeelement(qn("a:t"), {})
    new_t.text = new_text
    new_r.append(new_t)
    p_elem.append(new_r)


def _iter_text_frames(slide):
    """Yield every text_frame on the slide — both bare shapes AND table cells."""
    for shape in slide.shapes:
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
        elif shape.has_text_frame:
            yield shape.text_frame


def apply_edit(slide, find_text: str, replace_text: str) -> str:
    # Idempotency check.
    for tf in _iter_text_frames(slide):
        for para in tf.paragraphs:
            if replace_text in para.text:
                return "skip"

    # Pass 1: in-run replacement.
    for tf in _iter_text_frames(slide):
        for para in tf.paragraphs:
            for run in para.runs:
                if find_text in run.text:
                    run.text = run.text.replace(find_text, replace_text, 1)
                    return "ok"

    # Pass 2: merged-run fallback for substrings that cross runs.
    for tf in _iter_text_frames(slide):
        for para in tf.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if find_text in joined:
                _rewrite_paragraph_as_single_run(
                    para, joined.replace(find_text, replace_text, 1)
                )
                return "ok-merged"

    return "miss"


def main() -> None:
    prs = Presentation(SRC)
    print(f"Opened: {SRC}  ({len(prs.slides)} slides)")
    counts = {"ok": 0, "ok-merged": 0, "skip": 0, "miss": 0}
    for slide_1based, find, replace, label in EDITS:
        slide = prs.slides[slide_1based - 1]
        result = apply_edit(slide, find, replace)
        counts[result] += 1
        print(f"  [{result:>9}]  {label}")
    prs.save(SRC)
    print(f"\nsaved: {SRC}")
    print(f"counts: {counts}")


if __name__ == "__main__":
    main()
