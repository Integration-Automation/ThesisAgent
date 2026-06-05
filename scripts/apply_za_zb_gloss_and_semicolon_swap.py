"""Follow-up edits on exports/fang2026disentangling-zh-tw.pptx:

1. Gloss ``對抗 za`` and ``良性 zb`` at first use (slide 4).

   The previous clarity pass added Chinese glosses for ``VAE`` and prefixed
   the math operator ``互資訊`` for ``I(za;zb|Ep)``, but the variable names
   ``za`` / ``zb`` themselves were left bare. To a reader from an adjacent
   sub-field they read as "some letters" rather than "the two latent
   components the VAE encoder separates the prompt into". This pass adds
   ``對抗潛在向量(adversarial latent)`` / ``良性潛在向量(benign latent)``
   directly before the variable names on slide 4. Slide 16's recap already
   says ``對抗 latent za`` so the term is anchored, and slide 7 already
   uses ``za / zb 兩個 latent`` — both stay as-is (first-use rule = gloss
   the first chronological occurrence only).

2. Swap 4 prose ``;`` → ``,`` to match the new "use ， / , not ;" rule
   in ``CLAUDE.md`` "Content additions" section.

   Math notation ``;`` inside ``I(za;zb|Ep)`` is preserved on every slide
   (it's part of the operator syntax, not prose punctuation).

Idempotent: skips an edit whose replacement text is already present.
"""
from __future__ import annotations

import sys
from pathlib import Path

# Re-uses the in-run replace + merge-fallback helper from the
# clarity-rule script; nothing to import from there since the helper is
# small enough to copy.
_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from copy import deepcopy  # noqa: E402

from pptx import Presentation  # noqa: E402
from pptx.oxml.ns import qn  # noqa: E402

SRC = Path("exports/fang2026disentangling-zh-tw.pptx")


# (slide_1based, find_substring, replace_substring, short_label)
EDITS: list[tuple[int, str, str, str]] = [
    # ---- Slide 4 — combined: za/zb gloss + prose ; → , -------------------
    # The VAE bullet originally read:
    #   "...切成對抗 za 與良性 zb;訓練目標..."
    # New version glosses both latents AND swaps the prose ; to ,.
    (
        4,
        "切成對抗 za 與良性 zb;訓練目標最小化互資訊 I(za;zb|Ep)",
        "切成對抗潛在向量(adversarial latent)za 與良性潛在向量(benign latent)"
        "zb,訓練目標最小化互資訊 I(za;zb|Ep)",
        "S4: gloss za/zb at first use + swap prose ;",
    ),
    # ---- Slide 4 — Fiedler bullet's prose ; ------------------------------
    (
        4,
        "Fiedler 向量,圖 Laplacian 的第二小特徵向量;以及更高階特徵值",
        "Fiedler 向量,圖 Laplacian 的第二小特徵向量,以及更高階特徵值",
        "S4: prose ; → , in Fiedler clause",
    ),
    # ---- Slide 14 — ablation bullet's prose ; ----------------------------
    (
        14,
        "• 圖特徵貢獻 7.0 pp;單高階特徵值 2.3 pp",
        "• 圖特徵貢獻 7.0 pp,單高階特徵值 2.3 pp",
        "S14: prose ; → , between ablation rows",
    ),
    # ---- Slide 16 — contribution recap (mirror S4 VAE bullet's ; swap) ---
    (
        16,
        "與良性 latent zb;訓練目標最小化 I(za;zb|Ep)",
        "與良性 latent zb,訓練目標最小化 I(za;zb|Ep)",
        "S16: prose ; → , after za/zb pair",
    ),
]


def _rewrite_paragraph_as_single_run(para, new_text: str) -> None:
    """Replace a paragraph's runs with a single run carrying ``new_text``.

    Clones the first run's ``<a:rPr>`` so font, size, and (critically for
    the dark-mode contract) the explicit ``<a:solidFill>`` colour all
    survive the merge.
    """
    p_elem = para._p
    first_r = p_elem.find(qn("a:r"))
    rpr_clone = None
    if first_r is not None:
        rpr_elem = first_r.find(qn("a:rPr"))
        if rpr_elem is not None:
            rpr_clone = deepcopy(rpr_elem)
    for r in p_elem.findall(qn("a:r")):
        p_elem.remove(r)
    new_r = p_elem.makeelement(qn("a:r"), {})
    if rpr_clone is not None:
        new_r.append(rpr_clone)
    new_t = p_elem.makeelement(qn("a:t"), {})
    new_t.text = new_text
    new_r.append(new_t)
    p_elem.append(new_r)


def apply_edit(slide, find_text: str, replace_text: str) -> str:  # NOSONAR one-shot anchor-driven rewrite script; branchy by design (cf. ruff C901 scripts exemption)
    full_text = "\n".join(
        p.text for sh in slide.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    )
    if replace_text in full_text:
        return "skip"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if find_text in run.text:
                    run.text = run.text.replace(find_text, replace_text, 1)
                    return "ok"
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if find_text in joined:
                _rewrite_paragraph_as_single_run(
                    para, joined.replace(find_text, replace_text, 1)
                )
                return "ok-merged"
    return "miss"


def main() -> None:
    prs = Presentation(SRC)
    print(f"Opened: {SRC}  ({len(prs.slides)} slides)")
    counts = {"ok": 0, "ok-merged": 0, "skip": 0, "miss": 0}
    for slide_1based, find, replace, label in EDITS:
        slide = prs.slides[slide_1based - 1]
        result = apply_edit(slide, find, replace)
        counts[result] += 1
        print(f"  [{result:>9}]  {label}")
    prs.save(SRC)
    print(f"\nsaved: {SRC}")
    print(f"counts: {counts}")


if __name__ == "__main__":
    main()
