"""Convert exports/fang2026disentangling-zh-tw.pptx from light → dark mode.

The deck was generated 2026-05-19, before the dark-mode-default contract
landed; it ships in light mode despite carrying the bare (non-``-light``)
filename suffix the project's naming convention reserves for the dark
default. This script runs ``_apply_dark_mode(prs)`` in-place: solid-fills
every slide background with ``_DARK_SLIDE_BG`` (= ``#12151B``), walks
every shape / run / table cell and looks each RGB up in
``_LIGHT_TO_DARK_TEXT`` / ``_LIGHT_TO_DARK_FILL``, and ``_swap_text_colors``
promotes any leftover ``rgb is None`` or ``(0, 0, 0)`` run to ``#E5E7EB``
near-white — covering the 122 untouched body runs whose colour was
inherited from the light-mode theme default.

Filename stays the same (the bare suffix is the dark-default slot per
``CLAUDE.md`` "Read Subagents BEFORE Editing Any .pptx" table).
"""
from __future__ import annotations

import sys
from pathlib import Path

# Running this file directly puts `scripts/` on sys.path[0]; the project
# root needs to be there to import `autopapertoppt.*`.
_PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(_PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(_PROJECT_ROOT))

from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402

from autopapertoppt.exporters.pptx import _apply_dark_mode  # noqa: E402


SRC = Path("exports/fang2026disentangling-zh-tw.pptx")

# Per `.claude/agents/deck-design.md` "No red text" contract: any run currently
# painted `#C0392B` (the banned _BRAND_ACCENT red) must migrate to the
# sanctioned emphasis colour. In dark mode that is `#2DD4BF` (teal-400, the
# dark-mode equivalent of light-mode `_BRAND_HIGHLIGHT` teal-700). KPI values
# (the "the slide's punch line" role) keep their bold weight; the colour
# swap alone delivers the contract.
_BANNED_RED = RGBColor(0xC0, 0x39, 0x2B)
_SANCTIONED_DARK_TEAL = RGBColor(0x2D, 0xD4, 0xBF)


def _iter_text_runs(prs):
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            yield from para.runs
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                yield from para.runs


def _migrate_red_runs(prs) -> int:
    """Replace any run painted `#C0392B` with `#2DD4BF` (dark-mode teal)."""
    swapped = 0
    for run in _iter_text_runs(prs):
        try:
            rgb = run.font.color.rgb
        except Exception:
            continue
        if rgb == _BANNED_RED:
            run.font.color.rgb = _SANCTIONED_DARK_TEAL
            swapped += 1
    return swapped


def main() -> None:
    prs = Presentation(SRC)
    print(f"Opened: {SRC}  ({len(prs.slides)} slides)")
    _apply_dark_mode(prs)
    print("  dark-mode pass: background #12151B + light→dark text/fill swap")
    red_swapped = _migrate_red_runs(prs)
    print(f"  no-red-text contract: migrated {red_swapped} run(s) #C0392B → #2DD4BF")
    prs.save(SRC)
    print(f"saved: {SRC}")


if __name__ == "__main__":
    main()
