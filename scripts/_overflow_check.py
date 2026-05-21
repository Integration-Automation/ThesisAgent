"""Headless overflow inspection for one or more .pptx decks.

Walks every shape on every slide, estimates its rendered (wrapped) text
height with a per-font-size char-per-line heuristic, and flags shapes
that (a) overflow their declared box height or (b) extend past the
7.05" footer guard (where page numbers and slide footer live).

Usage:
    .venv\\Scripts\\python.exe -m scripts._overflow_check <pptx_path> [<pptx_path> ...]

Exit code 0 = all decks PASS, 1 = any violation found.
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu

# 7.05" in EMU — the body guard line. Body content (title / meta /
# subhead / body / kpi / paper_subtitle / rq_box) must not extend
# past this. The page_number and footer shapes are deliberately in
# the band beyond 7.05" — those are the footer territory itself.
FOOTER_GUARD_EMU = int(7.05 * 914400)
_FOOTER_BAND_SHAPES = frozenset({"page_number", "footer"})

# Approx chars per inch at given font sizes. Rough but matches the
# project's existing decks' wrap behaviour for default body text.
_CHARS_PER_INCH = {
    9: 14.0, 10: 12.5, 11: 11.5, 12: 10.5, 14: 9.0,
    16: 8.0, 18: 7.0, 20: 6.5, 24: 5.5, 28: 4.8, 30: 4.5, 36: 3.8,
}
_LINE_HEIGHT_FACTOR = 1.22  # line-height multiplier above raw font size


def _font_size_pt(run) -> int:
    sz = run.font.size
    if sz is None:
        return 12
    return max(8, sz.pt)


def _estimate_wrapped_height_emu(shape) -> int:
    tf = shape.text_frame
    width_in = (shape.width or Emu(0)) / 914400 or 5.0
    total_lines = 0.0
    weighted_lh_pt = 0.0
    for para in tf.paragraphs:
        text = "".join(r.text or "" for r in para.runs) or para.text or ""
        if not text:
            total_lines += 1
            weighted_lh_pt += 12.0 * _LINE_HEIGHT_FACTOR
            continue
        runs = list(para.runs)
        sz = _font_size_pt(runs[0]) if runs else 12
        cpi = _CHARS_PER_INCH.get(int(sz), 10.5)
        chars_per_line = max(1, int(cpi * width_in))
        text_lines = max(1, -(-len(text) // chars_per_line))  # ceil div
        total_lines += text_lines
        weighted_lh_pt += sz * _LINE_HEIGHT_FACTOR * text_lines
    avg_line_height_pt = weighted_lh_pt / max(1.0, total_lines)
    return int(avg_line_height_pt / 72.0 * 914400 * total_lines)


def _inspect(pptx_path: Path) -> list[tuple[int, str, str, int, int]]:
    """Returns a list of (slide_idx, shape_name, kind, rendered, limit)."""
    prs = Presentation(pptx_path)
    violations: list[tuple[int, str, str, int, int]] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            # Decorative shapes (`accent_top`, `accent_left`, etc.) have
            # an empty text_frame; estimating wrapped text on an empty
            # frame inflates to ~1 line-height which would false-flag
            # the 0.08" top accent bar. Skip when no actual text.
            if not (shape.text_frame.text or "").strip():
                continue
            name = shape.name or "?"
            top = shape.top or 0
            height = shape.height or 0
            auto = shape.text_frame.auto_size
            rendered = _estimate_wrapped_height_emu(shape)
            # TEXT_TO_FIT_SHAPE = PowerPoint will auto-shrink text to box;
            # SHAPE_TO_FIT_TEXT = PowerPoint will grow the box. Neither
            # produces a hard overflow at render time, so the height check
            # only applies when auto_size is NONE / disabled.
            strict = auto in (None, MSO_AUTO_SIZE.NONE)
            bottom = top + (rendered if strict else min(rendered, height))
            if strict and height and rendered > height:
                violations.append((idx, name, "overflows box", rendered, height))
            # page_number / footer shapes legitimately sit in the
            # footer band — don't flag them.
            if name in _FOOTER_BAND_SHAPES:
                continue
            if bottom > FOOTER_GUARD_EMU:
                violations.append((idx, name, "past footer guard", bottom, FOOTER_GUARD_EMU))
    return violations


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: python -m scripts._overflow_check <pptx> [<pptx> ...]")
        return 2
    bad = 0
    for arg in sys.argv[1:]:
        pptx_path = Path(arg)
        prs = Presentation(pptx_path)
        violations = _inspect(pptx_path)
        slides = len(list(prs.slides))
        shapes = sum(len(list(s.shapes)) for s in prs.slides)
        print(f"\noverflow check -- {pptx_path}")
        print(f"  slides: {slides}  shapes: {shapes}  violations: {len(violations)}")
        for idx, name, kind, rendered, limit in violations:
            ren_in = rendered / 914400
            lim_in = limit / 914400
            print(f"  slide {idx} shape {name!r}: {kind} -- {ren_in:.2f}\" vs {lim_in:.2f}\"")
        verdict = "PASS" if not violations else "FAIL"
        print(f"  verdict: {verdict}")
        if violations:
            bad += 1
    return 0 if bad == 0 else 1


if __name__ == "__main__":
    sys.exit(main())
