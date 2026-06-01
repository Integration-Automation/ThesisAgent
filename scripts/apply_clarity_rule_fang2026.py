"""Apply slide-deck-rules §8 (Content clarity & first-use context) to
exports/fang2026disentangling-zh-tw.pptx.

What's wrong with the deck pre-edit (audited 2026-05-28):

  Slide 3 uses ``ADA`` before defining it — the definition is buried on
  slide 9 inside the §4.2 evaluation block. ``Embedding clustering``
  appears without a one-clause function gloss.

  Slide 4 introduces ``VAE`` without expanding it to "Variational
  Autoencoder", and uses the math notation ``I(za;zb|Ep)`` without
  naming the operator (mutual information). ``DPI`` is parenthetically
  expanded but the gloss doesn't say WHAT the inequality guarantees.
  ``Fiedler 向量`` appears as if the reader already does spectral graph
  theory.

  Slide 5 KPI block uses ``EC`` shorthand in the comparator strings
  without spelling ``Embedding Clustering`` — this is the abbreviation's
  first standalone appearance.

Each edit is a single in-run text replacement. Run formatting (font.color
in particular — critical for the dark-mode contract) is preserved because
``run.text = new_text`` rewrites the text within the existing run and
doesn't touch the rPr / colour XML. Idempotent: skips an edit whose
replacement text is already present.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation


SRC = Path("exports/fang2026disentangling-zh-tw.pptx")


# (slide_1based, find_substring, replace_substring, short_label)
EDITS: list[tuple[int, str, str, str]] = [
    # ---- Slide 3 (research background + pain points) -----------------------
    (
        3,
        "規則式 ADA 在多樣 benchmark 上掉到 65.4%",
        "規則式 ADA(對抗偵測準確率,Adversarial Detection Accuracy)在多樣 benchmark 上掉到 65.4%",
        "S3: define ADA at first use",
    ),
    (
        3,
        "Embedding clustering:同一批 benchmark 只有 78.6%",
        "Embedding clustering(以 embedding 距離找異常的偵測法):同一批 benchmark 只有 78.6%",
        "S3: gloss Embedding clustering",
    ),

    # ---- Slide 4 (contributions) -------------------------------------------
    (
        4,
        "VAE 編碼器把 prompt 切成對抗 za 與良性 zb",
        "VAE(變分自編碼器,Variational Autoencoder)編碼器把 prompt 切成對抗 za 與良性 zb",
        "S4: define VAE at first use",
    ),
    (
        4,
        "訓練目標最小化 I(za;zb|Ep)",
        "訓練目標最小化互資訊 I(za;zb|Ep)",
        "S4: name the operator behind I(·;·|·)",
    ),
    (
        4,
        "透過 Data Processing Inequality(DPI)保證分離",
        "透過 Data Processing Inequality(DPI,限制條件互資訊在資料處理鏈中不會上升)保證分離",
        "S4: explain what DPI guarantees",
    ),
    (
        4,
        "在 za 的語意鄰域建圖,以譜分析(Fiedler 向量 + 高階特徵值)",
        "在 za 的語意鄰域建圖,以譜分析(Fiedler 向量,圖 Laplacian 的第二小特徵向量;以及更高階特徵值)",
        "S4: explain Fiedler vector",
    ),

    # ---- Slide 5 (KPI block) -----------------------------------------------
    (
        5,
        "對照組: 規則式 65.4 / EC 78.6 / AT 86.7",
        "對照組: 規則式 65.4 / EC(Embedding Clustering) 78.6 / AT(對抗訓練) 86.7",
        "S5: spell out EC + AT in comparator",
    ),
]


def apply_edit(slide, find_text: str, replace_text: str) -> str:
    """Find the first run containing ``find_text`` and replace it.

    Returns ``"ok"``, ``"skip"`` (replacement already present, idempotency),
    ``"miss"`` (substring not in any single run — most likely crosses runs
    and needs a follow-up to handle).
    """
    # Pass 1 — idempotency check on the slide as a whole.
    full_text = "\n".join(
        p.text for sh in slide.shapes
        if sh.has_text_frame
        for p in sh.text_frame.paragraphs
    )
    if replace_text in full_text:
        return "skip"

    # Pass 2 — locate the run that contains find_text and replace in-place.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if find_text in run.text:
                    run.text = run.text.replace(find_text, replace_text, 1)
                    return "ok"

    # Pass 3 — fallback: substring crosses runs. Concatenate all runs of
    # each paragraph, search, then rewrite the paragraph as a single run
    # cloning the first run's rPr.
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            joined = "".join(r.text for r in para.runs)
            if find_text in joined:
                new_text = joined.replace(find_text, replace_text, 1)
                _rewrite_paragraph_as_single_run(para, new_text)
                return "ok-merged"

    return "miss"


def _rewrite_paragraph_as_single_run(para, new_text: str) -> None:
    """Replace a paragraph's runs with a single run that has the new text.

    Clones the first run's rPr so font, colour, size carry over — critical
    for the dark-mode contract (run colour must remain a known palette RGB,
    not theme-default).
    """
    from copy import deepcopy
    from pptx.oxml.ns import qn

    p_elem = para._p
    first_r = p_elem.find(qn("a:r"))
    rpr_clone = None
    if first_r is not None:
        rpr_elem = first_r.find(qn("a:rPr"))
        if rpr_elem is not None:
            rpr_clone = deepcopy(rpr_elem)

    # Remove every existing <a:r> child.
    for r in p_elem.findall(qn("a:r")):
        p_elem.remove(r)

    # Build the replacement <a:r>.
    new_r = p_elem.makeelement(qn("a:r"), {})
    if rpr_clone is not None:
        new_r.append(rpr_clone)
    new_t = p_elem.makeelement(qn("a:t"), {})
    new_t.text = new_text
    new_r.append(new_t)
    p_elem.append(new_r)


def main() -> None:
    prs = Presentation(SRC)
    print(f"Opened: {SRC}  ({len(prs.slides)} slides)")

    counts = {"ok": 0, "ok-merged": 0, "skip": 0, "miss": 0}
    for slide_1based, find, replace, label in EDITS:
        slide = prs.slides[slide_1based - 1]
        result = apply_edit(slide, find, replace)
        counts[result] += 1
        print(f"  [{result:>9}]  {label}")

    prs.save(SRC)
    print(f"\nsaved: {SRC}")
    print(f"counts: {counts}")


if __name__ == "__main__":
    main()
