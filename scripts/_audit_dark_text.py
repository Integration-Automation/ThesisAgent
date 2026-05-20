"""Find dark-mode readability bugs in a rendered .pptx.

Two failure modes flagged:

A) "Black on dark slide bg" — text run whose colour is missing or too
   dark to read against the #12151B slide background:
   - ``run.font.color.rgb is None`` (inherits theme default → black)
   - ``rgb == (0,0,0)`` (explicit black)
   - ``rgb`` luminance below 60 AND not in the accepted text set

B) "White text inside white-fill box" — text whose run colour is
   light (luminance > 0.7 of 255) but the SHAPE FILL behind it is
   also light (luminance > 0.7) → text disappears into the box.
   Catches the `_RQ_BOX_FILL` class of bug where a near-white callout
   stays near-white after the dark-mode pass while the text it
   contains gets re-coloured to near-white.

Usage:
    .venv\\Scripts\\python.exe -m scripts._audit_dark_text <pptx_path>
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

# Colours we KNOW are intentional on a dark deck — the dark-mode text
# near-white, mid-greys, light-text on header fills. _BRAND_ACCENT
# (#C0392B warm red) is deliberately NOT in this set — red text was
# banned per the deck-design "No red text" contract.
_ACCEPTED_DARK_RUN_COLORS = {
    (0xE5, 0xE7, 0xEB),  # dark-mode body text
    (0x9C, 0xA3, 0xAF),  # dark-mode metadata grey
    (0x6B, 0x72, 0x80),  # dark-mode muted grey
    (0xFF, 0xFF, 0xFF),  # table-header white
}


def _luminance(rgb: tuple[int, int, int]) -> float:
    return 0.2126 * rgb[0] + 0.7152 * rgb[1] + 0.0722 * rgb[2]


def _iter_runs(prs):
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            yield from _shape_runs(slide_idx, shape, "")
            if shape.has_table:
                for r_idx, row in enumerate(shape.table.rows):
                    for c_idx, cell in enumerate(row.cells):
                        yield from _shape_runs(
                            slide_idx, cell, f"table[{r_idx},{c_idx}]"
                        )


def _shape_runs(slide_idx: int, shape_or_cell, where_hint: str):
    text_frame = getattr(shape_or_cell, "text_frame", None)
    if text_frame is None:
        return
    name = getattr(shape_or_cell, "name", "") or where_hint
    for p_idx, paragraph in enumerate(text_frame.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            try:
                rgb = run.font.color.rgb
            except (AttributeError, ValueError, TypeError):
                rgb = None
            text = (run.text or "")[:40]
            yield (slide_idx, name, p_idx, r_idx, rgb, text)


def _read_shape_fill_rgb(shape_or_cell) -> tuple[int, int, int] | None:
    fill = getattr(shape_or_cell, "fill", None)
    if fill is None:
        return None
    try:
        rgb = fill.fore_color.rgb
    except (AttributeError, ValueError, TypeError):
        return None
    if rgb is None:
        return None
    return (int(rgb[0]), int(rgb[1]), int(rgb[2]))


_LIGHT_LUMINANCE_THRESHOLD = 0.7 * 255  # > 178


def _check_contrast(prs, bad: list[str]) -> None:
    """Failure mode B: light text inside light-fill shape (= invisible)."""
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            fill_rgb = _read_shape_fill_rgb(shape)
            if fill_rgb is None or _luminance(fill_rgb) <= _LIGHT_LUMINANCE_THRESHOLD:
                continue
            _check_shape_text_against_light_fill(slide_idx, shape, fill_rgb, bad)


def _check_shape_text_against_light_fill(slide_idx, shape, fill_rgb, bad) -> None:
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    for p_idx, paragraph in enumerate(tf.paragraphs):
        for r_idx, run in enumerate(paragraph.runs):
            text = (run.text or "").strip()
            if not text:
                continue
            try:
                rgb = run.font.color.rgb
            except (AttributeError, ValueError, TypeError):
                continue
            if rgb is None:
                continue
            text_rgb = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            if _luminance(text_rgb) <= _LIGHT_LUMINANCE_THRESHOLD:
                continue
            bad.append(
                f"slide {slide_idx} {shape.name!r} p{p_idx}r{r_idx}  "
                f"LIGHT-ON-LIGHT  fill={fill_rgb} text-rgb={text_rgb}  "
                f"text={text[:40]!r}"
            )


def main(pptx_path: Path) -> int:
    prs = Presentation(pptx_path)
    bad: list[str] = []
    # Failure mode A — dark / unset text on dark slide bg.
    for slide_idx, name, p_idx, r_idx, rgb, text in _iter_runs(prs):
        if not text.strip():
            continue
        rgb_tuple = tuple(rgb) if rgb is not None else None
        if rgb_tuple is None:
            bad.append(
                f"slide {slide_idx} {name!r} p{p_idx}r{r_idx}  rgb=None  text={text!r}"
            )
            continue
        if rgb_tuple in _ACCEPTED_DARK_RUN_COLORS:
            continue
        if _luminance(rgb_tuple) < 80:
            bad.append(
                f"slide {slide_idx} {name!r} p{p_idx}r{r_idx}  "
                f"rgb={rgb_tuple}  lum={_luminance(rgb_tuple):.0f}  text={text!r}"
            )
    # Failure mode B — light text inside light-fill shape.
    _check_contrast(prs, bad)

    print(f"audit: {pptx_path.name}")
    print(f"runs flagged: {len(bad)}")
    for line in bad[:40]:
        print("  " + line)
    if len(bad) > 40:
        print(f"  ... ({len(bad) - 40} more)")
    return 0 if not bad else 1


if __name__ == "__main__":
    sys.exit(main(Path(sys.argv[1])))
