"""Dark-mode / contrast / red-text deck auditor (library home).

This is the importable home of the colour-contract audit. ``scripts/_audit_dark_text.py``
is a thin CLI wrapper around it, and ``thesisagents.exporters.review`` folds it
into the one-stop ``review_deck`` audit. The dark-mode / no-red / contrast
contracts (see ``.claude/agents/rules/deck-design.md``) are also pinned by the
``tests/test_exporters.py`` regression tests — but those only run on decks the
exporter *generates*. CLAUDE.md's "Read Subagents BEFORE Editing Any .pptx" rule
extends the same contracts to **hand-made decks**, which no test covers; this
module is what catches them.

Checks (each maps to a deck-design contract):

1. **Invisible run** — ``rgb is None`` or ``rgb == (0,0,0)``: inherits the theme
   colour and renders near-black on the dark slide background.
2. **Red text** — ``#C0392B`` (``_BRAND_ACCENT``): banned as a TEXT colour in
   both modes (reads as error, pattern-matches AI-generated emphasis).
3. **Light-on-light** — a near-white run inside a near-white-fill shape (both
   luminances > 0.7 × 255): the contrast-contract invisibility bug.
4. **Off-palette run (warning)** — a run whose colour is none of the sanctioned
   dark-mode run colours (``_ACCEPTED_DARK_RUN_COLORS``). Informational.

Importable: ``audit_deck(path) -> list[Issue]``.
"""
from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

# Sanctioned dark-mode run colours — the *values* of _LIGHT_TO_DARK_TEXT plus the
# near-white promotion target and the white table-header foreground (which sits on
# the dark navy header fill). Keep in sync with pptx.py if the palette changes.
_ACCEPTED_DARK_RUN_COLORS = frozenset({
    (0xE5, 0xE7, 0xEB),  # near-white body text
    (0x9C, 0xA3, 0xAF),  # mid grey
    (0x6B, 0x72, 0x80),  # muted grey
    (0x60, 0xA5, 0xFA),  # blue-400 highlight
    (0xFF, 0xFF, 0xFF),  # white table-header foreground (on navy fill)
})
_RED_ACCENT = (0xC0, 0x39, 0x2B)
_LIGHT_LUMA = 0.7 * 255  # luminance above which a colour counts as "light"


@dataclass(frozen=True)
class Issue:
    slide: int
    shape: str
    kind: str          # "invisible" | "red text" | "light-on-light" | "off-palette"
    detail: str
    hard: bool         # hard issues fail the deck; warnings do not


def _rgb_tuple(rgb):
    if rgb is None:
        return None
    return (int(rgb[0]), int(rgb[1]), int(rgb[2]))


def _luma(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    return 0.299 * r + 0.587 * g + 0.114 * b


def _shape_fill_rgb(shape):
    """The shape's solid fill colour as an (r,g,b) tuple, or None if it has no
    readable solid fill (background / pattern / inherited)."""
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color.type is not None:
            return _rgb_tuple(fill.fore_color.rgb)
    except (TypeError, ValueError, AttributeError):
        return None
    return None


def _iter_runs(shape):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                yield run


def audit_deck(path: str | Path) -> list[Issue]:
    """Return every dark-mode / contrast / red-text issue in the deck."""
    return audit_prs(Presentation(str(path)))


def audit_prs(prs) -> list[Issue]:
    """Audit an already-open Presentation. Split from ``audit_deck`` so the
    one-stop ``review_deck`` can load the file once and share the object."""
    issues: list[Issue] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            name = getattr(shape, "name", "?") or "?"
            fill_rgb = _shape_fill_rgb(shape)
            fill_light = fill_rgb is not None and _luma(fill_rgb) > _LIGHT_LUMA
            for run in _iter_runs(shape):
                rgb = _rgb_tuple(run.font.color.rgb if run.font.color and run.font.color.type
                                 else None)
                if rgb is None or rgb == (0, 0, 0):
                    issues.append(Issue(idx, name, "invisible",
                                        f"rgb={rgb} renders near-black on dark bg", True))
                    continue
                if rgb == _RED_ACCENT:
                    issues.append(Issue(idx, name, "red text", "#C0392B is banned", True))
                    continue
                if fill_light and _luma(rgb) > _LIGHT_LUMA:
                    issues.append(Issue(idx, name, "light-on-light",
                                        f"text {rgb} on fill {fill_rgb}", True))
                    continue
                if rgb not in _ACCEPTED_DARK_RUN_COLORS:
                    issues.append(Issue(idx, name, "off-palette",
                                        f"rgb={rgb} not a sanctioned dark-mode colour", False))
    return issues


def _report(path: str | Path) -> bool:
    issues = audit_deck(path)
    hard = [i for i in issues if i.hard]
    warn = [i for i in issues if not i.hard]
    print(f"dark-text audit — {path}")
    print(f"hard issues:   {len(hard)}")
    for i in hard:
        print(f"  slide {i.slide}, shape \"{i.shape}\": {i.kind} — {i.detail}")
    print(f"warnings:      {len(warn)}")
    for i in warn:
        print(f"  slide {i.slide}, shape \"{i.shape}\": {i.kind} — {i.detail}")
    verdict = "PASS" if not hard else "FAIL"
    print(f"verdict:       {verdict}")
    return not hard


def main(argv: list[str]) -> int:
    if not argv:
        print("usage: _audit_dark_text.py <deck.pptx> [more.pptx ...]")
        return 2
    failed = 0
    for i, path in enumerate(argv):
        if i:
            print()
        if not _report(path):
            failed += 1
    return failed


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
