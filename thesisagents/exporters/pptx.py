"""Full presentable PowerPoint deck.

Three rendering paths:

1. **Lightweight** (default, abstract only). Each paper gets cover · agenda ·
   overview · Background/Approach/Findings · references.

2. **Enriched-flat** (``PaperSummary`` with motivation/contributions/method/
   results/limitations/takeaways but no rich fields). Each paper's flat
   bullet sections become their own slides.

3. **Thesis-style** (``PaperSummary`` with rich fields: pain_points, research
   question, headline_metrics, technique_table, literature_table,
   method_sections, evaluation_sections, research_questions, rq_results,
   …). Each paper expands into a thesis-defence-style deck — multi-column
   pain-point quadrants, KPI callouts, technique-comparison tables, per-RQ
   result tables, contribution summary, limitations + future work, Q&A.

Deck-level features:

* 16:9 widescreen (13.33" × 7.5") matching modern slide tools' default.
* Page-number footer "N / total" on every slide except the cover.
* Section-title plus a horizontal rule beneath, matching the reference deck.

All template strings flow through ``i18n.py`` so the deck respects
``ExportOptions.language``. Shape names stay semantic (``title`` / ``meta``
/ ``body`` / ``footer`` / ``page_number``) so ``pptx_edit`` keeps working.
"""

from __future__ import annotations

import re
from collections.abc import Iterable
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from thesisagents.core.constants import EXPORT_PPTX
from thesisagents.core.exceptions import ExportError
from thesisagents.core.models import (
    ExportOptions,
    Paper,
    PaperCollection,
    PaperSummary,
    RqResult,
)
from thesisagents.exporters.base import Exporter
from thesisagents.exporters.i18n import strings_for, t

# ---------------------------------------------------------------------------
# Layout constants (16:9 widescreen)
# ---------------------------------------------------------------------------

_SLIDE_WIDTH = Inches(13.333)
_SLIDE_HEIGHT = Inches(7.5)

_MARGIN_X = Inches(0.5)
_BODY_WIDTH = Inches(12.333)            # SLIDE_WIDTH - 2*MARGIN
_TITLE_TOP = Inches(0.3)
_TITLE_HEIGHT = Inches(1.0)             # tall enough for a wrapped 2-line title at 30pt
_RULE_TOP = Inches(1.35)
_BODY_TOP = Inches(1.5)
_BODY_HEIGHT = Inches(5.4)
_FOOTER_Y = Inches(7.05)
_FOOTER_HEIGHT = Inches(0.3)
_PAGE_NUMBER_X = Inches(12.0)
_PAGE_NUMBER_WIDTH = Inches(1.2)

# Cover-slide specific (centered block)
_COVER_TITLE_TOP = Inches(1.6)
_COVER_TITLE_HEIGHT = Inches(1.6)
_COVER_SUBTITLE_TOP = Inches(3.3)
_COVER_SUBTITLE_HEIGHT = Inches(1.0)
_COVER_META_TOP = Inches(4.7)
_COVER_META_HEIGHT = Inches(1.6)

# Type sizes — calibrated so a 60-char Chinese bullet fits one wrapped line
# at body width 12.3" in 16:9, with comfortable line height.
_COVER_TITLE_PT = 44
_COVER_SUBTITLE_PT = 24
_COVER_META_PT = 18
_SECTION_TITLE_PT = 30
_SUBHEAD_PT = 22
_BODY_PT = 19
_TABLE_PT = 14
_FOOTER_PT = 10

# Kept for callers of ``_cap_bullets`` whose signature still accepts a
# (now-ignored) ``max_chars`` parameter. Slides rely on word_wrap rather
# than character-level truncation; see ``_clean`` / ``_sentences_to_bullets``.
_BULLET_MAX_CHARS = 60
_BULLET_MAX_CHARS_COL = 28

# Content-density caps to keep one slide from overflowing.
_BULLETS_PER_CELL_MAX = 6      # multi-column / quadrant cells; raised
                                # from 3 so authors don't lose substance
                                # to silent drops.
_METHOD_SECTIONS_PER_SLIDE = 2
_EVALUATION_SECTIONS_PER_SLIDE = 2
_REFERENCES_PER_SLIDE = 8       # split long bib lists across slides

# Colours (mirror the reference deck's palette)
_BRAND_DARK = RGBColor(0x1F, 0x3A, 0x66)
#: WARNING — DO NOT use _BRAND_ACCENT as a TEXT colour.
#: Red text in slide decks is consistently associated with errors,
#: warnings, and AI-generated KPI emphasis ("look at this number!").
#: The project bans red font runs entirely; use _BRAND_HIGHLIGHT (teal)
#: for emphasis instead. The constant is kept around in case a future
#: non-text accent shape (sparkline, badge, etc.) needs it, but every
#: existing TEXT callsite has been migrated to _BRAND_HIGHLIGHT.
#: See .claude/agents/rules/deck-design.md "No red text" contract.
_BRAND_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
#: Emphasis text colour — teal-700 (#0E7490). Replaces the banned red
#: _BRAND_ACCENT for KPI values, RQ question callouts, figure
#: captions, and other "this stands out" use cases. Pairs well with
#: bold; pairs cleanly with _BRAND_DARK navy as the secondary; reads
#: as professional/modern (think academic posters, not error banners).
#: Dark-mode pass swaps to teal-400 (#2DD4BF) via _LIGHT_TO_DARK_TEXT.
_BRAND_HIGHLIGHT = RGBColor(0x0E, 0x74, 0x90)
_BRAND_GREY = RGBColor(0x55, 0x55, 0x55)
_BRAND_LIGHT = RGBColor(0xAA, 0xAA, 0xAA)

# Per-language typography. (latin_family, east_asian_family). The Latin
# family also covers Cyrillic / Greek / Devanagari via Inter; the
# east-asian slot is what PowerPoint consults for CJK code points, so
# leaving it `None` would let PowerPoint pick a default that doesn't
# match the Latin choice. See ``deck-design`` agent doc for the
# full rationale. Inter degrades gracefully to Calibri on hosts without
# Inter installed.
_FONT_FAMILIES: dict[str, tuple[str, str | None]] = {
    "en":     ("Inter", None),
    "es":     ("Inter", None),
    "fr":     ("Inter", None),
    "de":     ("Inter", None),
    "pt":     ("Inter", None),
    "it":     ("Inter", None),
    "vi":     ("Inter", None),
    "id":     ("Inter", None),
    "ru":     ("Inter", None),
    "hi":     ("Inter", "Nirmala UI"),
    "zh-tw":  ("Inter", "Microsoft JhengHei UI"),
    "zh-cn":  ("Inter", "Microsoft YaHei UI"),
    "ja":     ("Inter", "Yu Gothic UI"),
    "ko":     ("Inter", "Malgun Gothic"),
}
_DEFAULT_FONT_FAMILY: tuple[str, str | None] = ("Inter", None)

# Accent geometry (set on every content slide by the typography /
# accent pass so a stock blank layout still reads as a designed deck).
_ACCENT_TOP_HEIGHT = Inches(0.08)
_ACCENT_LEFT_WIDTH = Inches(0.4)

# Dark-mode palette (post-build recolour, opt-in via
# ``ExportOptions.dark_mode``).
_DARK_SLIDE_BG = RGBColor(0x12, 0x15, 0x1B)

# Light-palette RGB → dark-palette RGB mapping for TEXT colours. Keys
# are 3-tuples (R, G, B) since python-pptx's RGBColor is tuple-comparable
# but we want to match by raw int components.
_LIGHT_TO_DARK_TEXT: dict[tuple[int, int, int], tuple[int, int, int]] = {
    (0x1F, 0x3A, 0x66): (0xE5, 0xE7, 0xEB),  # _BRAND_DARK      → near-white text
    (0x55, 0x55, 0x55): (0x9C, 0xA3, 0xAF),  # _BRAND_GREY      → mid grey
    (0xAA, 0xAA, 0xAA): (0x6B, 0x72, 0x80),  # _BRAND_LIGHT     → muted grey
    (0x0E, 0x74, 0x90): (0x2D, 0xD4, 0xBF),  # _BRAND_HIGHLIGHT → bright teal-400
    # _BRAND_ACCENT (#C0392B) intentionally NOT mapped — red text was
    # banned per the deck-design "No red text" contract, and the
    # `test_pptx_no_red_text_runs` regression test fails if any run
    # ever writes that colour. If a run shows up with it the test
    # catches it BEFORE we reach this swap layer.
}

# Light-palette RGB → dark-palette RGB mapping for SHAPE / CELL FILLS
# and cell-border lines. Keeps the navy header on tables but lightens
# its tone slightly so it reads against the dark slide background.
_LIGHT_TO_DARK_FILL: dict[tuple[int, int, int], tuple[int, int, int]] = {
    # _BRAND_DARK accent bars + accent_left + table header fill
    (0x1F, 0x3A, 0x66): (0x3B, 0x5A, 0xA0),
    # _TABLE_ROW_ALT → dark row stripe
    (0xF4, 0xF6, 0xF9): (0x1F, 0x23, 0x2C),
    # Pure white table rows → near-black
    (0xFF, 0xFF, 0xFF): (0x16, 0x1A, 0x22),
    # _TABLE_DIVIDER → muted grey-blue rule
    (0xD0, 0xD7, 0xE2): (0x3D, 0x44, 0x52),
    # _RQ_BOX_FILL (research-question callout box) → dark navy tint.
    # Without this swap the box stays near-white while the text inside
    # is re-coloured to near-white = white-on-white = invisible. This
    # specific bug is what the dark-mode contrast contract guards.
    (0xF3, 0xF6, 0xFA): (0x1E, 0x26, 0x38),
}
_BRAND_RULE = RGBColor(0xCC, 0xCC, 0xCC)
_RQ_BOX_FILL = RGBColor(0xF3, 0xF6, 0xFA)
_RQ_BOX_BORDER = RGBColor(0x1F, 0x3A, 0x66)
_TABLE_HEADER_FILL = RGBColor(0x1F, 0x3A, 0x66)
_TABLE_HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
_TABLE_ROW_ALT = RGBColor(0xF4, 0xF6, 0xF9)
_TABLE_DIVIDER = RGBColor(0xD0, 0xD7, 0xE2)  # row divider — soft grey-blue
_TABLE_HEADER_RULE = RGBColor(0x1F, 0x3A, 0x66)  # heavy nav rule under header

# ---------------------------------------------------------------------------
# Abstract segmentation (fallback when summary is absent / lightweight only)
# ---------------------------------------------------------------------------

_LABEL_RE = re.compile(
    r"\b(BACKGROUND|INTRODUCTION|MOTIVATION|OBJECTIVE[S]?|AIM[S]?|"
    r"METHOD[S]?|APPROACH|DESIGN|MATERIALS?|"
    r"RESULT[S]?|FINDINGS?|EVALUATION|"
    r"CONCLUSION[S]?|DISCUSSION|TAKEAWAYS?|IMPLICATIONS?|"
    r"LIMITATIONS?|FUTURE WORK)\s*:\s*",
    re.IGNORECASE,
)
_SENTENCE_RE = re.compile(r"(?<=[.!?])\s+(?=[A-Z(])")
_HEADER_KEYS = {
    "background": ("background", "introduction", "motivation"),
    "approach": (
        "objective", "objectives", "aim", "aims",
        "method", "methods", "approach", "design", "material", "materials",
    ),
    "results": (
        "result", "results", "finding", "findings", "evaluation",
        "conclusion", "conclusions", "discussion", "takeaway", "takeaways",
        "implication", "implications", "limitation", "limitations",
        "future work",
    ),
}


def _segment_abstract(abstract: str) -> dict[str, list[str]]:
    cleaned = " ".join((abstract or "").split())
    if not cleaned:
        return {"background": [], "approach": [], "results": []}
    labelled = _labelled_split(cleaned)
    if labelled is not None:
        return labelled
    return _equal_third_split(cleaned)


def _labelled_split(text: str) -> dict[str, list[str]] | None:
    matches = list(_LABEL_RE.finditer(text))
    if len(matches) < 2:
        return None
    buckets: dict[str, list[str]] = {"background": [], "approach": [], "results": []}
    for i, match in enumerate(matches):
        label = match.group(1).lower()
        start = match.end()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        body = text[start:end].strip().rstrip(".")
        if not body:
            continue
        sentences = [s.strip() for s in _SENTENCE_RE.split(body) if s.strip()]
        target = _label_to_bucket(label)
        buckets[target].extend(sentences)
    if not any(buckets.values()):
        return None
    return buckets


def _label_to_bucket(label: str) -> str:
    for bucket_name, keys in _HEADER_KEYS.items():
        if label in keys:
            return bucket_name
    return "results"


def _equal_third_split(text: str) -> dict[str, list[str]]:
    sentences = [s.strip() for s in _SENTENCE_RE.split(text) if s.strip()]
    if not sentences:
        return {"background": [text], "approach": [], "results": []}
    n = len(sentences)
    if n <= 2:
        return {"background": sentences, "approach": [], "results": []}
    third = max(1, n // 3)
    return {
        "background": sentences[:third],
        "approach": sentences[third : 2 * third],
        "results": sentences[2 * third :],
    }


# ---------------------------------------------------------------------------
# Build context
# ---------------------------------------------------------------------------


class _BuildContext:
    __slots__ = ("language", "include_abstract", "_page_total")

    def __init__(self, language: str, include_abstract: bool) -> None:
        self.language = language
        self.include_abstract = include_abstract
        self._page_total = 0

    def set_page_total(self, total: int) -> None:
        self._page_total = total

    @property
    def page_total(self) -> int:
        return self._page_total


# ---------------------------------------------------------------------------
# Exporter
# ---------------------------------------------------------------------------


class PptxExporter(Exporter):
    format = EXPORT_PPTX
    extension = "pptx"

    def export(self, collection: PaperCollection, options: ExportOptions) -> Path:
        try:
            presentation = self._build(collection, options)
        except Exception as err:
            raise ExportError(self.format, f"render failed: {err}") from err
        out_path = self.resolve_out_path(collection, options)
        presentation.save(str(out_path))
        return out_path

    def _build(
        self, collection: PaperCollection, options: ExportOptions
    ) -> Presentation:
        prs = (
            Presentation(options.pptx_template)
            if options.pptx_template
            else Presentation()
        )
        prs.slide_width = _SLIDE_WIDTH
        prs.slide_height = _SLIDE_HEIGHT
        blank = prs.slide_layouts[6]
        total = len(collection)
        ctx = _BuildContext(language=options.language, include_abstract=options.include_abstract)
        _add_cover_slide(prs, blank, collection, ctx)
        if total > 1:
            _add_agenda_slide(prs, blank, collection, ctx)
        for index, paper in enumerate(collection.papers, start=1):
            _add_paper_slides(prs, blank, index, total, paper, ctx)
        if total > 0:
            _add_references_slide(prs, blank, collection, ctx)
        # Drop lower-priority slides if a positive cap is configured.
        # ``None`` and any non-positive int both disable the cap so the
        # full rich-tier deck flows through unchanged.
        budget = options.max_slides_per_paper
        if budget is not None and budget > 0:
            _trim_to_slide_budget(
                prs,
                language=ctx.language,
                paper_count=max(total, 1),
                budget_per_paper=budget,
            )
        # Page numbers are stamped AFTER trim so they reflect the final total.
        _stamp_page_numbers(prs, ctx.language)
        # Visual identity passes — applied last so they affect every shape
        # placed by every builder (including page numbers). See the
        # ``deck-design`` subagent doc for rationale.
        _apply_typography(prs, ctx.language)
        _decorate_with_accents(prs)
        if options.dark_mode:
            _apply_dark_mode(prs)
        return prs


# ---------------------------------------------------------------------------
# Per-paper dispatch
# ---------------------------------------------------------------------------


def _add_paper_slides(
    prs: Presentation, layout, index: int, total: int, paper: Paper, ctx: _BuildContext
) -> None:
    if total > 1:
        _add_section_divider(prs, layout, index, total, paper, ctx)
    _add_overview_slide(prs, layout, index, total, paper, ctx)
    if not ctx.include_abstract:
        return
    summary = paper.summary
    if summary is not None and summary.has_rich_fields():
        _add_rich_summary_slides(prs, layout, index, total, paper, summary, ctx)
        return
    if summary is not None and not summary.is_empty():
        _add_flat_summary_slides(prs, layout, index, total, paper, summary, ctx)
        return
    if paper.abstract:
        _add_abstract_split_slides(prs, layout, index, total, paper, ctx)


#: Slide plan entries: (predicate(summary) -> bool, builder).
#: Used by _add_rich_summary_slides; split into two halves around per-RQ
#: result slides so RQ-intro → per-RQ tables → contribution summary stays
#: in narrative order.
def _build_rich_plan(summary: PaperSummary):
    return [
        (bool(summary.pain_points), _add_pain_points_slide),
        (
            bool(summary.research_question) and not summary.pain_points,
            _add_research_question_slide,
        ),
        (
            bool(summary.contributions_detailed or summary.headline_metrics),
            _add_contributions_detailed_slide,
        ),
        (bool(summary.technique_table), _add_technique_table_slide),
        (bool(summary.literature_table), _add_literature_table_slide),
        (bool(summary.system_flow), _add_system_overview_slide),
        (bool(summary.figures), _add_figure_slides),
        (bool(summary.method_sections), _add_method_details_slides),
        (bool(summary.evaluation_sections), _add_evaluation_slide),
        (bool(summary.paper_tables), _add_paper_table_slides),
        (bool(summary.research_questions), _add_research_questions_slide),
    ]


def _build_rich_plan_tail(summary: PaperSummary):
    return [
        (
            bool(summary.contributions_detailed or summary.core_observation),
            _add_contribution_summary_slide,
        ),
        (bool(summary.limitations or summary.future_work), _add_limitations_future_slide),
    ]


def _add_rich_summary_slides(
    prs: Presentation, layout, index: int, total: int,
    paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    _ = index, total  # reserved for future per-slide footers
    for predicate, builder in _build_rich_plan(summary):
        if predicate:
            builder(prs, layout, paper, summary, ctx)
    rq_lookup = {rq_id: question for rq_id, question in summary.research_questions}
    for rq in summary.rq_results:
        _add_rq_result_slide(prs, layout, paper, rq, ctx, rq_lookup=rq_lookup)
    for predicate, builder in _build_rich_plan_tail(summary):
        if predicate:
            builder(prs, layout, paper, summary, ctx)
    _add_qa_slide(prs, layout, paper, ctx)


def _add_flat_summary_slides(
    prs: Presentation, layout, index: int, total: int,
    paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    sections: tuple[tuple[str, tuple[str, ...]], ...] = (
        ("section_motivation", summary.motivation),
        ("section_contributions", summary.contributions),
        ("section_method", summary.method),
        ("section_results", summary.results),
        ("section_limitations", summary.limitations),
        ("section_takeaways", summary.takeaways),
    )
    for key, bullets in sections:
        if bullets:
            _add_content_slide(prs, layout, index, total, paper, key, list(bullets), ctx)


def _add_abstract_split_slides(
    prs: Presentation, layout, index: int, total: int,
    paper: Paper, ctx: _BuildContext,
) -> None:
    buckets = _segment_abstract(paper.abstract)
    _add_content_slide(
        prs, layout, index, total, paper,
        "section_background", buckets["background"], ctx,
    )
    if buckets["approach"]:
        _add_content_slide(
            prs, layout, index, total, paper,
            "section_approach", buckets["approach"], ctx,
        )
    _add_content_slide(
        prs, layout, index, total, paper,
        "section_findings",
        buckets["results"] or buckets["background"], ctx,
    )


# ---------------------------------------------------------------------------
# Slide builders — cover / agenda / divider / overview / refs / qa
# ---------------------------------------------------------------------------


def _add_cover_slide(
    prs: Presentation, layout, collection: PaperCollection, ctx: _BuildContext
) -> None:
    slide = prs.slides.add_slide(layout)
    title_text = _cover_title(collection, ctx)
    _add_textbox(
        slide, name="title", text=title_text,
        left=_MARGIN_X, top=_COVER_TITLE_TOP,
        width=_BODY_WIDTH, height=_COVER_TITLE_HEIGHT,
        font_pt=_COVER_TITLE_PT, bold=True, colour=_BRAND_DARK,
        align=PP_ALIGN.CENTER,
        shrink_to_fit=True,
    )
    if len(collection) == 1:
        subtitle = _english_subtitle(collection.papers[0])
        # Skip the subtitle when it would duplicate the title verbatim — this
        # is the common case for a monolingual English deck. The subtitle is
        # only useful when the cover title was localised and we want the
        # original paper title surfaced underneath.
        if subtitle and subtitle.strip() != title_text.strip():
            _add_textbox(
                slide, name="subtitle", text=subtitle,
                left=_MARGIN_X, top=_COVER_SUBTITLE_TOP,
                width=_BODY_WIDTH, height=_COVER_SUBTITLE_HEIGHT,
                font_pt=_COVER_SUBTITLE_PT, colour=_BRAND_GREY,
                align=PP_ALIGN.CENTER,
            )
    meta_text = _cover_subtitle(collection, ctx)
    _add_textbox(
        slide, name="meta", text=meta_text,
        left=_MARGIN_X, top=_COVER_META_TOP,
        width=_BODY_WIDTH, height=_COVER_META_HEIGHT,
        font_pt=_COVER_META_PT, colour=_BRAND_GREY,
        align=PP_ALIGN.CENTER,
    )


def _add_agenda_slide(
    prs: Presentation, layout, collection: PaperCollection, ctx: _BuildContext
) -> None:
    slide = _new_section_slide(prs, layout, t(ctx.language, "agenda"))
    bullets = [_agenda_line(i + 1, p, ctx) for i, p in enumerate(collection.papers)]
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=_BODY_TOP,
        width=_BODY_WIDTH, height=_BODY_HEIGHT,
        font_pt=_BODY_PT,
    )


def _add_section_divider(
    prs: Presentation, layout, index: int, total: int, paper: Paper, ctx: _BuildContext
) -> None:
    slide = prs.slides.add_slide(layout)
    _add_textbox(
        slide, name="title",
        text=t(ctx.language, "paper_n_of_m", n=index, m=total),
        left=_MARGIN_X, top=Inches(2.6),
        width=_BODY_WIDTH, height=Inches(0.8),
        font_pt=_FOOTER_PT + 8, colour=_BRAND_GREY,
        align=PP_ALIGN.CENTER,
        shrink_to_fit=True,
    )
    _add_textbox(
        slide, name="meta",
        text=_clean(paper.title),
        left=_MARGIN_X, top=Inches(3.4),
        width=_BODY_WIDTH, height=Inches(2.4),
        font_pt=_SECTION_TITLE_PT + 2, bold=True, colour=_BRAND_DARK,
        align=PP_ALIGN.CENTER,
        shrink_to_fit=True,
    )


def _add_overview_slide(
    prs: Presentation, layout, index: int, total: int, paper: Paper, ctx: _BuildContext
) -> None:
    slide = _new_section_slide(
        prs, layout,
        _clean(paper.title),
        font_pt=_SECTION_TITLE_PT,
    )
    meta_text = _overview_meta(paper, ctx)
    _add_textbox(
        slide, name="meta", text=meta_text,
        left=_MARGIN_X, top=_BODY_TOP,
        width=_BODY_WIDTH, height=Inches(1.1),
        font_pt=_BODY_PT - 1, colour=_BRAND_GREY,
    )
    bullets = list(_overview_bullets(paper, ctx))
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=_BODY_TOP + Inches(1.2),
        width=_BODY_WIDTH, height=Inches(4.3),
        font_pt=_BODY_PT,
    )


def _add_references_slide(
    prs: Presentation, layout, collection: PaperCollection, ctx: _BuildContext
) -> None:
    slide = _new_section_slide(prs, layout, t(ctx.language, "references"))
    bullets = [_reference_line(i + 1, p, ctx) for i, p in enumerate(collection.papers)]
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=_BODY_TOP,
        width=_BODY_WIDTH, height=_BODY_HEIGHT,
        font_pt=_BODY_PT - 2,
    )
    _add_footer(slide, t(ctx.language, "footer_references"))


def _add_qa_slide(prs: Presentation, layout, paper: Paper, ctx: _BuildContext) -> None:
    slide = prs.slides.add_slide(layout)
    _add_textbox(
        slide, name="title", text=t(ctx.language, "section_qa"),
        left=_MARGIN_X, top=Inches(2.6),
        width=_BODY_WIDTH, height=Inches(1.2),
        font_pt=_COVER_TITLE_PT + 4, bold=True, colour=_BRAND_DARK,
        align=PP_ALIGN.CENTER,
        shrink_to_fit=True,
    )
    _add_textbox(
        slide, name="meta", text=t(ctx.language, "footer_qa"),
        left=_MARGIN_X, top=Inches(4.2),
        width=_BODY_WIDTH, height=Inches(0.8),
        font_pt=_COVER_SUBTITLE_PT, colour=_BRAND_GREY,
        align=PP_ALIGN.CENTER,
    )
    _ = paper  # unused but kept for symmetry with other slide builders


# ---------------------------------------------------------------------------
# Rich-tier slide builders
# ---------------------------------------------------------------------------


_PAIN_POINTS_PER_SLIDE = 4   # 2 columns × 2 rows


def _add_pain_points_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    """Pain-points quadrant slide, paginated when more than 4 sections
    are supplied. The research-question callout sits on the first slide
    only — subsequent pages are full quadrants of pain points."""
    title = t(ctx.language, "section_pain_points")
    sections = list(summary.pain_points)
    if not sections:
        return
    chunks = [
        sections[i : i + _PAIN_POINTS_PER_SLIDE]
        for i in range(0, len(sections), _PAIN_POINTS_PER_SLIDE)
    ]
    for chunk_index, chunk in enumerate(chunks):
        chunk_title = title
        if len(chunks) > 1:
            chunk_title = f"{title} ({chunk_index + 1}/{len(chunks)})"
        slide = _new_section_slide(prs, layout, chunk_title)
        _add_paper_subtitle(slide, paper, ctx)
        _render_multi_column(
            slide,
            sections=tuple(chunk),
            left=_MARGIN_X, top=Inches(1.7),
            width=_BODY_WIDTH, height=Inches(4.4),
            columns=2,
        )
        if chunk_index == 0 and summary.research_question:
            _add_rq_callout(
                slide, summary.research_question,
                left=_MARGIN_X, top=Inches(6.2),
                width=_BODY_WIDTH, height=Inches(0.8),
            )


def _add_research_question_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    # If we already rendered the RQ inline on the pain-points slide, skip the
    # dedicated slide. Otherwise (no pain_points but a standalone RQ), give it
    # its own slide.
    if summary.pain_points:
        return
    title = t(ctx.language, "section_research_question")
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    _add_rq_callout(
        slide, summary.research_question,
        left=_MARGIN_X, top=Inches(2.5),
        width=_BODY_WIDTH, height=Inches(2.0),
    )


_MAX_STACKS_PER_SLIDE = 5   # standalone stacks slide
_FOOTER_GUARD = Inches(7.0)


def _add_contributions_detailed_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    """Render contributions and KPI metrics on dedicated slides.

    We always split: even if the paper only has 2 contributions, the KPI
    metrics get their own slide so the layout is predictable and no slide
    has to balance "stacks + KPI" inside a fixed height.
    """
    if summary.contributions_detailed:
        _add_stacks_slide(
            prs, layout, paper, ctx,
            title=t(ctx.language, "section_contributions"),
            stacks=summary.contributions_detailed,
        )
    if summary.headline_metrics:
        _add_kpi_slide(prs, layout, paper, ctx, summary.headline_metrics)


def _add_stacks_slide(
    prs: Presentation, layout, paper: Paper, ctx: _BuildContext,
    *, title: str, stacks,
) -> None:
    """Render stacked sections. Paginates when stacks exceed
    ``_MAX_STACKS_PER_SLIDE`` so author bullets are never silently
    dropped — instead the title gets ``(1/N)`` and overflow spills onto
    the next slide.
    """
    stack_list = list(stacks)
    if not stack_list:
        return
    chunks = [
        stack_list[i : i + _MAX_STACKS_PER_SLIDE]
        for i in range(0, len(stack_list), _MAX_STACKS_PER_SLIDE)
    ]
    for chunk_index, chunk in enumerate(chunks):
        chunk_title = title
        if len(chunks) > 1:
            chunk_title = f"{title} ({chunk_index + 1}/{len(chunks)})"
        slide = _new_section_slide(prs, layout, chunk_title)
        _add_paper_subtitle(slide, paper, ctx)
        cursor = Inches(1.7)
        for heading, body in chunk:
            cursor = _add_stacked_section(slide, heading, body, cursor)


_KPI_PER_SLIDE = 6


def _add_kpi_slide(
    prs: Presentation, layout, paper: Paper, ctx: _BuildContext, metrics,
) -> None:
    """Paginate KPI metrics in chunks of ``_KPI_PER_SLIDE`` so an author
    who supplies 10 metrics gets two slides, not one with 6 and four
    silently dropped."""
    metric_list = list(metrics)
    if not metric_list:
        return
    chunks = [
        metric_list[i : i + _KPI_PER_SLIDE]
        for i in range(0, len(metric_list), _KPI_PER_SLIDE)
    ]
    for chunk_index, chunk in enumerate(chunks):
        title = t(ctx.language, "label_headline_metrics")
        if len(chunks) > 1:
            title = f"{title} ({chunk_index + 1}/{len(chunks)})"
        slide = _new_section_slide(prs, layout, title)
        _add_paper_subtitle(slide, paper, ctx)
        kpi_top = Inches(1.7)
        kpi_height = Inches(min(0.55 * len(chunk), 5.0))
        _add_kpi_lines(
            slide, tuple(chunk), ctx,
            left=_MARGIN_X, top=kpi_top,
            width=_BODY_WIDTH, height=kpi_height,
        )


def _add_technique_table_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    title = t(ctx.language, "section_technique_overview")
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    rows = [(t(ctx.language, "label_technique") if False else "Technique",
             t(ctx.language, "label_role") if False else "Role")]
    # Use the table header as supplied by the language pack key fallbacks.
    # For now we just label per language.
    if ctx.language.startswith("zh"):
        rows = [("技術", "於本研究之角色")]
    elif ctx.language == "ja":
        rows = [("技術", "本研究での役割")]
    else:
        rows = [("Technique", "Role in this paper")]
    rows.extend(summary.technique_table)
    _add_table(
        slide, rows=rows,
        left=_MARGIN_X, top=Inches(1.7),
        width=_BODY_WIDTH, height=Inches(5.0),
        col_widths=(Inches(3.8), Inches(8.5)),
    )


def _add_literature_table_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    title = t(ctx.language, "section_literature_positioning")
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    rows = summary.literature_table
    if not rows:
        return
    cols = len(rows[0])
    col_widths = _equal_col_widths(_BODY_WIDTH, cols)
    _add_table(
        slide, rows=rows,
        left=_MARGIN_X, top=Inches(1.7),
        width=_BODY_WIDTH, height=Inches(5.0),
        col_widths=col_widths,
    )


def _add_system_overview_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    title = t(ctx.language, "section_system_overview")
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    bullets = [f"{i + 1}. {step}" for i, step in enumerate(summary.system_flow)]
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=Inches(1.7),
        width=_BODY_WIDTH, height=Inches(5.0),
        font_pt=_BODY_PT,
    )


def _add_method_details_slides(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    sections = summary.method_sections
    chunk_size = _METHOD_SECTIONS_PER_SLIDE
    chunks = [sections[i : i + chunk_size] for i in range(0, len(sections), chunk_size)]
    for chunk_index, chunk in enumerate(chunks):
        title = t(ctx.language, "section_method_details")
        if len(chunks) > 1:
            title = f"{title} ({chunk_index + 1}/{len(chunks)})"
        slide = _new_section_slide(prs, layout, title)
        _add_paper_subtitle(slide, paper, ctx)
        cursor = Inches(1.7)
        for heading, bullets in chunk:
            cursor = _add_subsection(
                slide, heading, bullets, cursor, width=_BODY_WIDTH,
            )


def _add_evaluation_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    sections = summary.evaluation_sections
    chunk_size = _EVALUATION_SECTIONS_PER_SLIDE
    chunks = [sections[i : i + chunk_size] for i in range(0, len(sections), chunk_size)]
    for chunk_index, chunk in enumerate(chunks):
        title = t(ctx.language, "section_evaluation")
        if len(chunks) > 1:
            title = f"{title} ({chunk_index + 1}/{len(chunks)})"
        slide = _new_section_slide(prs, layout, title)
        _add_paper_subtitle(slide, paper, ctx)
        _render_multi_column(
            slide,
            sections=tuple(chunk),
            left=_MARGIN_X, top=Inches(1.7),
            width=_BODY_WIDTH, height=Inches(5.0),
            columns=min(2, len(chunk)),
        )


def _add_figure_slides(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    """One slide per figure: image at top, caption below, description bullets.

    The image is loaded from ``figure.image_path`` and scaled to fit a
    fixed area, preserving aspect ratio. A missing / unreadable file is
    rendered as a placeholder bullet so the deck still flows.
    """
    base_title = t(ctx.language, "section_figure")
    for index, (caption, image_path, description) in enumerate(summary.figures, 1):
        title = (
            f"{base_title} {index}/{len(summary.figures)}"
            if len(summary.figures) > 1 else base_title
        )
        slide = _new_section_slide(prs, layout, title)
        _add_paper_subtitle(slide, paper, ctx)
        _add_figure_image(
            slide, image_path,
            left=_MARGIN_X, top=Inches(1.7),
            max_width=_BODY_WIDTH, max_height=Inches(3.6),
        )
        # Caption sits between image and bullets, full-width, italic-ish weight.
        _add_textbox(
            slide, name="subhead",
            text=f"{t(ctx.language, 'label_caption')}: {_clean(caption)}",
            left=_MARGIN_X, top=Inches(5.4),
            width=_BODY_WIDTH, height=Inches(0.5),
            font_pt=_SUBHEAD_PT - 4, bold=True, colour=_BRAND_GREY,
            shrink_to_fit=True,
        )
        if description:
            _add_bullet_box(
                slide, name="body",
                bullets=_cap_bullets(description, max_count=4),
                left=_MARGIN_X, top=Inches(5.95),
                width=_BODY_WIDTH, height=Inches(1.05),
                font_pt=_BODY_PT - 2,
            )


def _add_figure_image(
    slide, image_path: str, *, left, top, max_width, max_height,
) -> None:
    from pathlib import Path as PathT

    path = PathT(image_path)
    if not path.is_file():
        _add_textbox(
            slide, name="body",
            text=f"[figure unavailable: {path.name}]",
            left=left, top=top, width=max_width, height=Inches(0.5),
            # Muted grey for a placeholder/error state — not a headline.
            # Was red, then briefly navy; settled on grey because this
            # surface is contextual chrome, not "this stands out".
            font_pt=_BODY_PT, colour=_BRAND_GREY,
        )
        return
    # python-pptx scales by aspect ratio if we pass only height (or
    # width), but we don't know the image's intrinsic size up-front.
    # Strategy: insert with max_height; if the resulting width overflows
    # the body, re-insert with max_width instead. PIL would let us
    # compute aspect ratio first but we avoid the extra dep.
    pic = slide.shapes.add_picture(
        str(path), left, top, height=max_height,
    )
    if pic.width > max_width:
        # Drop + reinsert constrained by width.
        sp = pic._element  # noqa: SLF001  # python-pptx needs raw XML access
        sp.getparent().remove(sp)
        pic = slide.shapes.add_picture(
            str(path), left, top, width=max_width,
        )
    # Centre horizontally if narrower than body
    if pic.width < max_width:
        pic.left = left + Emu((max_width - pic.width) // 2)
    pic.name = "figure"


def _add_paper_table_slides(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    """One slide per paper-table: render the table plus analysis bullets."""
    base_title = t(ctx.language, "section_paper_table")
    for index, (caption, rows, analysis) in enumerate(summary.paper_tables, 1):
        if not rows:
            continue
        title = (
            f"{base_title} {index}/{len(summary.paper_tables)}"
            if len(summary.paper_tables) > 1 else base_title
        )
        slide = _new_section_slide(prs, layout, title)
        _add_paper_subtitle(slide, paper, ctx)
        _add_textbox(
            slide, name="subhead",
            text=f"{t(ctx.language, 'label_caption')}: {_clean(caption)}",
            left=_MARGIN_X, top=Inches(1.65),
            width=_BODY_WIDTH, height=Inches(0.55),
            # Mid grey reads as a caption label (matches the figure-slide
            # caption style at line ~887) rather than competing with the
            # paper-table itself for the eye. Was red, then briefly navy.
            font_pt=_SUBHEAD_PT - 2, bold=True, colour=_BRAND_GREY,
            shrink_to_fit=True,
        )
        cols = len(rows[0])
        col_widths = _equal_col_widths(_BODY_WIDTH, cols)
        _add_table(
            slide, rows=rows,
            left=_MARGIN_X, top=Inches(2.3),
            width=_BODY_WIDTH, height=Inches(3.0),
            col_widths=col_widths,
        )
        if analysis:
            _add_bullet_box(
                slide, name="body",
                bullets=_cap_bullets(analysis, max_count=6),
                left=_MARGIN_X, top=Inches(5.6),
                width=_BODY_WIDTH, height=Inches(1.4),
                font_pt=_BODY_PT,
            )


def _add_research_questions_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    title = t(ctx.language, "section_research_questions")
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    bullets = [f"{rq_id}: {question}" for rq_id, question in summary.research_questions]
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=Inches(1.8),
        width=_BODY_WIDTH, height=Inches(4.9),
        font_pt=_BODY_PT + 1,
    )


def _add_rq_result_slide(
    prs: Presentation, layout, paper: Paper, rq: RqResult, ctx: _BuildContext,
    *, rq_lookup: dict[str, str] | None = None,
) -> None:
    title = f"{t(ctx.language, 'section_results_for')} {rq.rq_id}"
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    # Prefer the verbatim question from ``research_questions`` so the RQ
    # slide carries the paper's actual wording rather than a short label
    # the author may have used in the RqResult.question field. The lookup
    # is keyed by rq_id; falls back to whatever the RqResult provided.
    question_text = rq.question
    if rq_lookup:
        full_question = rq_lookup.get(rq.rq_id, "")
        if full_question:
            question_text = full_question
    _add_textbox(
        slide, name="rq_question", text=question_text,
        left=_MARGIN_X, top=Inches(1.65),
        width=_BODY_WIDTH, height=Inches(0.55),
        # Teal highlight — this is the actual research question being
        # answered on this slide; the eye should land on it before the
        # results table below. Was red, then briefly navy; teal carries
        # the "thoughtful, intentional" tone without the warning vibe.
        font_pt=_SUBHEAD_PT - 2, bold=True, colour=_BRAND_HIGHLIGHT,
        shrink_to_fit=True,
    )
    if rq.table:
        cols = len(rq.table[0])
        col_widths = _equal_col_widths(_BODY_WIDTH, cols)
        _add_table(
            slide, rows=rq.table,
            left=_MARGIN_X, top=Inches(2.3),
            width=_BODY_WIDTH, height=Inches(3.0),
            col_widths=col_widths,
        )
    if rq.analysis:
        analysis_top = Inches(5.6) if rq.table else Inches(2.4)
        # Show up to 6 analysis bullets so authors don't lose argument
        # detail to a silent 3-bullet cap.
        _add_bullet_box(
            slide, name="body", bullets=_cap_bullets(rq.analysis, max_count=6),
            left=_MARGIN_X, top=analysis_top,
            width=_BODY_WIDTH, height=Inches(1.4 if rq.table else 4.5),
            font_pt=_BODY_PT,
        )


def _add_contribution_summary_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    """Always split contribution-summary stacks vs core_observation callout."""
    if summary.contributions_detailed:
        _add_stacks_slide(
            prs, layout, paper, ctx,
            title=t(ctx.language, "section_contribution_summary"),
            stacks=summary.contributions_detailed,
        )
    if summary.core_observation:
        slide = _new_section_slide(
            prs, layout, t(ctx.language, "label_core_observation")
        )
        _add_paper_subtitle(slide, paper, ctx)
        _add_rq_callout(
            slide, summary.core_observation,
            left=_MARGIN_X, top=Inches(2.5),
            width=_BODY_WIDTH, height=Inches(2.0),
        )


def _add_limitations_future_slide(
    prs: Presentation, layout, paper: Paper, summary: PaperSummary, ctx: _BuildContext,
) -> None:
    title = (
        f"{t(ctx.language, 'section_limitations')} & "
        f"{t(ctx.language, 'section_future_work')}"
    )
    slide = _new_section_slide(prs, layout, title)
    _add_paper_subtitle(slide, paper, ctx)
    sections: list[tuple[str, tuple[str, ...]]] = []
    if summary.limitations:
        sections.append(
            (t(ctx.language, "section_limitations"), tuple(summary.limitations))
        )
    if summary.future_work:
        sections.append(
            (t(ctx.language, "section_future_work"), tuple(summary.future_work))
        )
    _render_multi_column(
        slide,
        sections=tuple(sections),
        left=_MARGIN_X, top=Inches(1.7),
        width=_BODY_WIDTH, height=Inches(5.0),
        columns=2,
    )


# ---------------------------------------------------------------------------
# Legacy content slide (used by flat / abstract fallback paths)
# ---------------------------------------------------------------------------


def _add_content_slide(
    prs: Presentation, layout, index: int, total: int,
    paper: Paper, section_key: str, sentences: list[str], ctx: _BuildContext,
) -> None:
    section_title = t(ctx.language, section_key)
    slide = _new_section_slide(prs, layout, section_title)
    _add_textbox(
        slide, name="meta",
        text=(
            f"{_clean(paper.title)} · "
            f"{t(ctx.language, 'paper_n_of_m', n=index, m=total)}"
        ),
        left=_MARGIN_X, top=_BODY_TOP,
        width=_BODY_WIDTH, height=Inches(0.5),
        font_pt=_BODY_PT - 4, colour=_BRAND_GREY,
    )
    fallback = [_clean(paper.abstract)]
    raw_bullets = _sentences_to_bullets(sentences) or fallback
    bullets = _cap_bullets(raw_bullets, max_count=6, max_chars=_BULLET_MAX_CHARS)
    _add_bullet_box(
        slide, name="body", bullets=bullets,
        left=_MARGIN_X, top=_BODY_TOP + Inches(0.6),
        width=_BODY_WIDTH, height=Inches(4.7),
        font_pt=_BODY_PT,
    )


# ---------------------------------------------------------------------------
# Text helpers
# ---------------------------------------------------------------------------


def _cover_title(collection: PaperCollection, ctx: _BuildContext) -> str:
    if len(collection) == 1:
        return _clean(collection.papers[0].title)
    prefix = t(ctx.language, "paper_review_prefix")
    return f"{prefix} {collection.query.keywords}"


def _english_subtitle(paper: Paper) -> str:
    # For single-paper decks, surface the paper's title verbatim under the
    # (possibly translated) cover title. When the title is already in English
    # this duplicates — that's fine for now and easier than language-detecting.
    return paper.title


def _cover_subtitle(collection: PaperCollection, ctx: _BuildContext) -> str:
    paper_count = len(collection)
    source_list = ", ".join(collection.query.sources)
    if paper_count == 1:
        paper = collection.papers[0]
        authors = ", ".join(paper.authors[:3])
        if len(paper.authors) > 3:
            authors += " et al."
        year = paper.year or t(ctx.language, "year_no_date")
        venue = paper.venue or paper.source
        return f"{authors}\n{year} · {venue}"
    span_bits: list[str] = []
    if collection.query.year_from or collection.query.year_to:
        span_bits.append(
            f"{collection.query.year_from or '—'} – {collection.query.year_to or '—'}"
        )
    span = f" ({'; '.join(span_bits)})" if span_bits else ""
    count_text = _paper_count_text(paper_count, ctx)
    return f"{count_text} · {source_list}{span}"


def _paper_count_text(count: int, ctx: _BuildContext) -> str:
    key = "paper_count_singular" if count == 1 else "paper_count_plural"
    return t(ctx.language, key, n=count)


def _agenda_line(index: int, paper: Paper, ctx: _BuildContext) -> str:
    title = _clean(paper.title)
    year = paper.year or t(ctx.language, "year_no_date")
    return f"{index:>2}. {title} ({year})"


def _reference_line(index: int, paper: Paper, ctx: _BuildContext) -> str:
    authors_short = paper.authors[0] if paper.authors else "Anon"
    if len(paper.authors) > 1:
        authors_short += " et al."
    year = paper.year or t(ctx.language, "year_no_date")
    ident = paper.doi or paper.arxiv_id or paper.url
    title = _clean(paper.title)
    return f"{index:>2}. {authors_short} ({year}). {title} — {ident}"


def _overview_meta(paper: Paper, ctx: _BuildContext) -> str:
    authors_line = (
        ", ".join(paper.authors) if paper.authors else t(ctx.language, "no_authors_dash")
    )
    year_part = str(paper.year) if paper.year else t(ctx.language, "year_no_date")
    venue_part = f" · {paper.venue}" if paper.venue else ""
    return (
        f"{authors_line}\n{year_part}{venue_part}\n"
        f"{t(ctx.language, 'label_source')}: {paper.source}"
    )


def _overview_bullets(paper: Paper, ctx: _BuildContext) -> Iterable[str]:
    s = strings_for(ctx.language)
    if paper.url:
        yield f"{s['label_url']}: {paper.url}"
    if paper.doi:
        yield f"{s['label_doi']}: {paper.doi}"
    if paper.arxiv_id:
        yield f"{s['label_arxiv']}: {paper.arxiv_id}"
    if paper.pdf_url:
        yield f"{s['label_pdf']}: {paper.pdf_url}"
    if paper.citation_count is not None:
        yield f"{s['label_citations']}: {paper.citation_count}"
    yield f"{s['label_bibtex_key']}: {paper.bibtex_key()}"


def _sentences_to_bullets(sentences: list[str]) -> list[str]:
    """Whitespace-clean each sentence; do not truncate.

    Why no per-bullet ellipsis: a deck whose bullets read like
    "蒐集到 40…" is worse than a deck whose bullets wrap to two lines.
    Word wrap is on for every body box, so authors stay responsible for
    bullet length and the layout never silently mangles a sentence.
    """
    bullets: list[str] = []
    for sentence in sentences:
        stripped = " ".join(sentence.split()).rstrip(".")
        if stripped:
            bullets.append(stripped)
    return bullets


def _cap_bullets(
    bullets,
    max_count: int = _BULLETS_PER_CELL_MAX,
    max_chars: int = _BULLET_MAX_CHARS,  # kept for back-compat; ignored
) -> list[str]:
    """Cap the *number* of bullets shown; never truncate a bullet's text.

    A trailing "(+N more)" marker calls out genuine overflow by count so
    authors notice they exceeded the cap, without silently chewing
    characters off a kept bullet.
    """
    del max_chars  # intentionally ignored; see docstring
    sliced = [" ".join(b.split()) for b in bullets[:max_count]]
    overflow = len(bullets) - len(sliced)
    if overflow > 0:
        sliced.append(f"(+{overflow} more)")
    return sliced


def _clean(text: str) -> str:
    """Collapse runs of whitespace; never truncate.

    Why no truncation: a slide deck whose titles and subtitles end in
    ``…`` hides information from the reader without warning. We've left
    word-wrap on for every textbox, so long titles can wrap to a second
    line cleanly. Authors are responsible for keeping titles short
    enough to look good; the renderer never silently lops them off.
    """
    return " ".join((text or "").split())


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------


def _new_section_slide(
    prs: Presentation, layout, title: str, *, font_pt: int = _SECTION_TITLE_PT,
):
    slide = prs.slides.add_slide(layout)
    # ``shrink_to_fit`` lets a long title (e.g. a verbatim paper title) wrap
    # within the fixed-height title box and PowerPoint scales the font down
    # so the text never crosses the horizontal rule below.
    _add_textbox(
        slide, name="title", text=_clean(title),
        left=_MARGIN_X, top=_TITLE_TOP,
        width=_BODY_WIDTH, height=_TITLE_HEIGHT,
        font_pt=font_pt, bold=True, colour=_BRAND_DARK,
        shrink_to_fit=True,
    )
    _add_horizontal_rule(slide, top=_RULE_TOP)
    return slide


def _add_paper_subtitle(slide, paper: Paper, ctx: _BuildContext) -> None:
    _ = ctx  # reserved for future per-language formatting
    # The right-hand element is the publication venue, NOT the fetcher.
    # ``paper.source`` is internal plumbing (arxiv / openalex / pubmed) and
    # only useful as a fallback when the source didn't fill ``venue`` in.
    publication = paper.venue or paper.source
    text = f"{_clean(paper.title)}  ·  {publication}"
    _add_textbox(
        slide, name="paper_subtitle", text=text,
        left=_MARGIN_X, top=Inches(1.4),
        width=_BODY_WIDTH, height=Inches(0.28),
        font_pt=_BODY_PT - 5, colour=_BRAND_GREY,
        shrink_to_fit=True,
    )


def _add_horizontal_rule(slide, *, top) -> None:
    from pptx.shapes.connector import Connector  # noqa: F401  # not used directly

    line = slide.shapes.add_connector(
        connector_type=1,  # straight
        begin_x=_MARGIN_X, begin_y=top,
        end_x=Emu(_SLIDE_WIDTH - _MARGIN_X), end_y=top,
    )
    line.line.color.rgb = _BRAND_RULE
    line.line.width = Pt(0.75)


def _add_textbox(
    slide, *, name: str, text: str, left, top, width, height,
    font_pt: int, bold: bool = False, colour: RGBColor | None = None,
    align: PP_ALIGN | None = None,
    shrink_to_fit: bool = False,
) -> None:
    """Render a textbox.

    ``shrink_to_fit`` flips on PowerPoint's "Shrink text on overflow"
    auto-size (``MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE``). Use it on the
    title-area boxes so a long title that we now decline to truncate still
    gets rendered inside its allotted height — PowerPoint shrinks the font
    at open time rather than letting the text bleed past the horizontal
    rule.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    text_frame = box.text_frame
    text_frame.word_wrap = True
    if shrink_to_fit:
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        if align is not None:
            paragraph.alignment = align
        for run in paragraph.runs:
            run.font.size = Pt(font_pt)
            run.font.bold = bold
            if colour is not None:
                run.font.color.rgb = colour


def _add_bullet_box(
    slide, *, name: str, bullets: list[str], left, top, width, height, font_pt: int,
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = name
    text_frame = box.text_frame
    text_frame.word_wrap = True
    if not bullets:
        text_frame.text = ""
        return
    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = Pt(font_pt)
            # ALWAYS set the run colour explicitly. A run with
            # ``font.color.rgb = None`` inherits the theme's body-text
            # colour (which renders as black) and the dark-mode
            # post-pass cannot swap it because there's no source RGB
            # to look up in the mapping. See deck-design.md
            # "Dark-mode contract" — every text-adding helper sets a
            # palette colour, no exceptions.
            run.font.color.rgb = _BRAND_DARK


def _add_footer(slide, text: str) -> None:
    _add_textbox(
        slide, name="footer", text=text,
        left=_MARGIN_X, top=_FOOTER_Y,
        width=_BODY_WIDTH - Inches(1.5), height=_FOOTER_HEIGHT,
        font_pt=_FOOTER_PT, colour=_BRAND_LIGHT,
    )


def _add_stacked_section(slide, heading: str, body: str, cursor) -> int:
    """Render an inline (heading bold + body grey) block; return next cursor Y.

    Body height is sized for a 2-line wrap at the current body font so a
    full-sentence contribution doesn't visually spill into the next
    subhead. Body text is not truncated — see ``_sentences_to_bullets``
    for the rationale.
    """
    head_height = Inches(0.42)
    body_height = Inches(0.85)
    gap = Inches(0.05)
    _add_textbox(
        slide, name="subhead", text=_clean(heading),
        left=_MARGIN_X, top=cursor,
        width=_BODY_WIDTH, height=head_height,
        font_pt=_SUBHEAD_PT, bold=True, colour=_BRAND_DARK,
    )
    _add_textbox(
        slide, name="body", text=" ".join((body or "").split()),
        left=Inches(0.7), top=cursor + head_height,
        width=_BODY_WIDTH - Inches(0.2), height=body_height,
        font_pt=_BODY_PT, colour=_BRAND_GREY,
    )
    return cursor + head_height + body_height + gap


def _add_subsection(slide, heading: str, bullets, cursor, *, width) -> int:
    head_height = Inches(0.5)
    # Subsection bullets — show up to 6 so authors don't lose detail.
    # Method/eval slides paginate at the section-list level, so 6 per
    # subsection won't push past the footer in the worst case.
    capped = _cap_bullets(bullets, max_count=6)
    bullet_height = Inches(0.5 * max(1, len(capped)))
    _add_textbox(
        slide, name="subhead", text=_clean(heading),
        left=_MARGIN_X, top=cursor,
        width=width, height=head_height,
        font_pt=_SUBHEAD_PT, bold=True, colour=_BRAND_DARK,
    )
    _add_bullet_box(
        slide, name="body", bullets=capped,
        left=Inches(0.7), top=cursor + head_height,
        width=width - Inches(0.2), height=bullet_height,
        font_pt=_BODY_PT,
    )
    return cursor + head_height + bullet_height + Inches(0.25)


def _render_multi_column(
    slide, *, sections, left, top, width, height, columns: int = 2,
) -> None:
    if not sections:
        return
    cols = columns
    col_w = (width - Inches(0.2 * (cols - 1))) / cols
    rows = (len(sections) + cols - 1) // cols
    if rows == 0:
        return
    row_h = height / rows
    for idx, (heading, bullets) in enumerate(sections):
        r, c = divmod(idx, cols)
        x = left + c * (col_w + Inches(0.2))
        y = top + r * row_h
        _add_textbox(
            slide, name="subhead", text=_clean(heading),
            left=x, top=y,
            width=col_w, height=Inches(0.5),
            font_pt=_SUBHEAD_PT, bold=True, colour=_BRAND_DARK,
        )
        _add_bullet_box(
            slide, name="body",
            bullets=_cap_bullets(bullets, max_chars=_BULLET_MAX_CHARS_COL),
            left=x, top=y + Inches(0.55),
            width=col_w, height=row_h - Inches(0.65),
            font_pt=_BODY_PT,
        )


def _add_rq_callout(slide, text: str, *, left, top, width, height) -> None:
    """A boxed highlight: filled rectangle + bold text on top."""
    from pptx.enum.shapes import MSO_SHAPE

    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.name = "rq_box"
    rect.fill.solid()
    rect.fill.fore_color.rgb = _RQ_BOX_FILL
    rect.line.color.rgb = _RQ_BOX_BORDER
    rect.line.width = Pt(1.0)
    rect.text_frame.word_wrap = True
    rect.text_frame.margin_left = Inches(0.2)
    rect.text_frame.margin_right = Inches(0.2)
    rect.text_frame.margin_top = Inches(0.1)
    rect.text_frame.margin_bottom = Inches(0.1)
    rect.text_frame.text = text
    for paragraph in rect.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(_BODY_PT)
            run.font.bold = True
            run.font.color.rgb = _BRAND_DARK


def _add_kpi_lines(
    slide, metrics, ctx: _BuildContext, *, left, top, width, height,
) -> None:
    """Each metric → one paragraph in a textbox: `Label: VALUE (baseline X)`."""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.name = "kpi"
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.text = ""
    baseline_label = t(ctx.language, "label_baseline")
    for index, (label, value, baseline) in enumerate(metrics):
        paragraph = (
            text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        )
        paragraph.text = ""
        run_label = paragraph.add_run()
        run_label.text = f"• {label}: "
        run_label.font.size = Pt(_BODY_PT)
        run_label.font.color.rgb = _BRAND_GREY
        run_value = paragraph.add_run()
        run_value.text = str(value)
        run_value.font.size = Pt(_BODY_PT + 2)
        run_value.font.bold = True
        # Teal accent for KPI numbers — they're the slide's punch line
        # (a 2.3x speedup, a 78% F1, etc.). Bold + teal makes them pop
        # without using red, which would read as error/warning. Was red,
        # then briefly navy; teal restores a real emphasis colour.
        # See deck-design.md "No red text" contract.
        run_value.font.color.rgb = _BRAND_HIGHLIGHT
        if baseline:
            run_base = paragraph.add_run()
            run_base.text = f"   ({baseline_label}: {baseline})"
            run_base.font.size = Pt(_BODY_PT - 2)
            run_base.font.color.rgb = _BRAND_LIGHT


def _add_table(
    slide, *, rows, left, top, width, height, col_widths,
) -> None:
    """Render a clean academic-style table.

    Styling rules (mirror published thesis-defence decks, not the default
    PowerPoint table look):

    * No default black grid lines — every cell border is set to noFill
      first, then specific rules are added back where they help readability.
    * Header row: navy fill, white bold text, with a thick (1.5pt) navy
      bottom rule below it for emphasis (the rule sits in the data row's
      top edge, not the header's bottom, so it doesn't double up).
    * Data rows: alternate very-light-blue / white background; thin
      (0.5pt) grey-blue rule between adjacent data rows.
    * Cell vertical alignment: middle, so short labels and longer
      descriptions in the same row sit on a shared baseline.
    * First column of body rows: bold, slightly emphasised — most tables
      in this project use the leftmost cell as a row label.
    """
    if not rows:
        return
    row_count = len(rows)
    col_count = len(rows[0])
    table_shape = slide.shapes.add_table(
        rows=row_count, cols=col_count, left=left, top=top, width=width, height=height,
    )
    table = table_shape.table
    for col_index, w in enumerate(col_widths):
        table.columns[col_index].width = w
    for r, row_values in enumerate(rows):
        for c, value in enumerate(row_values):
            _style_table_cell(table.cell(r, c), str(value), r, c)


def _style_table_cell(cell, value: str, r: int, c: int) -> None:
    """Apply academic-style formatting to one cell.

    Split out from ``_add_table`` so the cognitive-complexity budget
    fits — borders + fills + font + alignment all live here.
    """
    cell.text = value
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    text_frame = cell.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(_TABLE_PT)
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = _TABLE_HEADER_FG
            else:
                run.font.color.rgb = _BRAND_DARK
                if c == 0:
                    # Row-label column gets a slightly heavier weight.
                    run.font.bold = True
    _set_cell_fill(cell, r)
    _clear_cell_borders(cell)
    if r == 1:
        # Heavy rule below the header — drawn as the data row's TOP
        # border so the visual width adds up cleanly.
        _set_cell_border(cell, "T", Pt(1.5), _TABLE_HEADER_RULE)
    elif r > 1:
        # Thin separator between adjacent data rows.
        _set_cell_border(cell, "T", Pt(0.5), _TABLE_DIVIDER)


def _set_cell_fill(cell, r: int) -> None:
    cell.fill.solid()
    if r == 0:
        cell.fill.fore_color.rgb = _TABLE_HEADER_FILL
    elif r % 2 == 0:
        cell.fill.fore_color.rgb = _TABLE_ROW_ALT
    else:
        cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)


def _clear_cell_borders(cell) -> None:
    """Set every edge of ``cell`` to noFill so the table style's default
    grid lines disappear. The XML structure for a cell's border is
    ``<a:tcPr>/<a:lnX>/<a:noFill/>`` where X ∈ {L, R, T, B}.
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    for edge in ("L", "R", "T", "B"):
        tag = qn(f"a:ln{edge}")
        existing = tc_pr.find(tag)
        if existing is not None:
            tc_pr.remove(existing)
        ln = tc_pr.makeelement(tag, {"w": "0"}, nsmap=None)
        ln.append(ln.makeelement(qn("a:noFill"), {}, nsmap=None))
        tc_pr.append(ln)


def _set_cell_border(cell, edge: str, width, colour: RGBColor) -> None:
    """Add a solid-fill border on one edge of ``cell``.

    ``edge`` must be one of L / R / T / B; ``width`` is a ``pptx.util.Pt``
    (or any EMU-aware value). Replaces an existing border on that edge.
    """
    tc_pr = cell._tc.get_or_add_tcPr()
    tag = qn(f"a:ln{edge}")
    existing = tc_pr.find(tag)
    if existing is not None:
        tc_pr.remove(existing)
    ln = tc_pr.makeelement(
        tag,
        {"w": str(int(width)), "cap": "flat", "cmpd": "sng", "algn": "ctr"},
        nsmap=None,
    )
    solid = ln.makeelement(qn("a:solidFill"), {}, nsmap=None)
    rgb_hex = f"{colour[0]:02X}{colour[1]:02X}{colour[2]:02X}"
    solid.append(solid.makeelement(qn("a:srgbClr"), {"val": rgb_hex}, nsmap=None))
    ln.append(solid)
    ln.append(ln.makeelement(qn("a:prstDash"), {"val": "solid"}, nsmap=None))
    ln.append(ln.makeelement(qn("a:round"), {}, nsmap=None))
    tc_pr.append(ln)


def _equal_col_widths(total: Emu, cols: int) -> tuple[Emu, ...]:
    each = total // cols
    return tuple(each for _ in range(cols))


#: Slide-category priority for the ``max_slides_per_paper`` trim.
#: Higher = keep; lower = drop first when over budget. Categories not
#: listed here default to ``_DEFAULT_PRIORITY``.
_CATEGORY_PRIORITY: dict[str, int] = {
    "cover": 100,
    "references": 100,
    "contributions": 95,
    "metrics": 95,
    "overview": 92,
    "research_question": 90,
    "core_observation": 90,
    "pain_points": 85,
    "research_questions": 80,
    "agenda": 78,
    "rq_results": 70,
    "limitations_future": 70,
    "method_details": 60,
    "evaluation": 60,
    "system_overview": 55,
    "technique_table": 55,
    "figure": 50,
    "paper_table": 50,
    "qa": 45,
    "contribution_summary": 40,
    "literature_table": 40,
    "section_divider": 38,
}
_DEFAULT_PRIORITY = 65


def _trim_to_slide_budget(
    prs: Presentation, *, language: str, paper_count: int, budget_per_paper: int,
) -> None:
    """Drop the lowest-priority slides until the per-paper budget fits.

    For multi-paper decks the budget scales linearly with paper count
    (so ``--max-slides 12`` on a 3-paper deck keeps up to 36 slides
    plus the shared cover / agenda / references). The cover, agenda,
    references and core-deliverable slides are always preferred.
    """
    total_budget = max(1, budget_per_paper * paper_count)
    current = len(prs.slides)
    if current <= total_budget:
        return
    categories = _categorise_slides(prs, language)
    indexed = sorted(
        range(current),
        key=lambda i: (
            _CATEGORY_PRIORITY.get(categories[i], _DEFAULT_PRIORITY),
            -i,  # break ties by dropping later slides first
        ),
    )
    drop_count = current - total_budget
    drop_set = set(indexed[:drop_count])
    import logging

    logging.getLogger("thesisagents.exporters.pptx").info(
        "max_slides trim: budget=%d, current=%d → dropping %d "
        "slides (categories=%s)",
        total_budget, current, drop_count,
        sorted({categories[i] for i in drop_set}),
    )
    sld_id_list = prs.slides._sldIdLst  # noqa: SLF001  # python-pptx exposes only via _ attr
    elements = list(sld_id_list)
    for i in sorted(drop_set, reverse=True):
        sld_id_list.remove(elements[i])


def _categorise_slides(prs: Presentation, language: str) -> list[str]:
    """Tag each slide with a category so the trim can drop by priority.

    Title-based heuristic: each slide's ``title`` shape text is matched
    against i18n strings to recover the category. Slide 0 is always the
    cover (paper title varies so we can't match it).
    """
    strings = strings_for(language)
    label_to_cat = _build_label_to_category_map(strings)
    markers = {
        "limitations": strings.get("section_limitations", ""),
        "future_work": strings.get("section_future_work", ""),
        "results_for": strings.get("section_results_for", ""),
        "paper_n_of_m_prefix": strings["paper_n_of_m"].split("{", 1)[0].strip(),
    }
    out: list[str] = []
    for i, slide in enumerate(prs.slides):
        if i == 0:
            out.append("cover")
            continue
        title = _slide_title_text(slide)
        category = _classify_title(title, label_to_cat, markers)
        # Single-paper decks emit an "overview" slide right after the cover
        # whose title is the paper title (so it doesn't match any i18n key).
        if category == "unknown" and i == 1:
            category = "overview"
        out.append(category)
    return out


def _build_label_to_category_map(strings: dict[str, str]) -> dict[str, str]:
    return {
        strings["agenda"]: "agenda",
        strings["references"]: "references",
        strings["section_pain_points"]: "pain_points",
        strings["section_research_question"]: "research_question",
        strings["section_contributions"]: "contributions",
        strings["label_headline_metrics"]: "metrics",
        strings["section_technique_overview"]: "technique_table",
        strings["section_literature_positioning"]: "literature_table",
        strings["section_system_overview"]: "system_overview",
        strings["section_method_details"]: "method_details",
        strings["section_evaluation"]: "evaluation",
        strings["section_research_questions"]: "research_questions",
        strings["section_contribution_summary"]: "contribution_summary",
        strings["label_core_observation"]: "core_observation",
        strings["section_qa"]: "qa",
        strings["section_figure"]: "figure",
        strings["section_paper_table"]: "paper_table",
    }


def _slide_title_text(slide) -> str:
    for shape in slide.shapes:
        if shape.has_text_frame and shape.name == "title":
            return shape.text_frame.text.strip()
    return ""


def _classify_title(
    title: str, label_to_cat: dict[str, str], markers: dict[str, str],
) -> str:
    limitations = markers["limitations"]
    future_work = markers["future_work"]
    if limitations and future_work and limitations in title and future_work in title:
        return "limitations_future"
    if markers["results_for"] and title.startswith(markers["results_for"]):
        return "rq_results"
    if (
        markers["paper_n_of_m_prefix"]
        and title.startswith(markers["paper_n_of_m_prefix"])
    ):
        return "section_divider"
    # Longest matching label wins ("Research Questions" beats "Research Question").
    category = "unknown"
    best_match_len = 0
    for label, name in label_to_cat.items():
        if label and title.startswith(label) and len(label) > best_match_len:
            category = name
            best_match_len = len(label)
    return category


def _stamp_page_numbers(prs: Presentation, language: str) -> None:
    _ = language  # reserved for future locale-specific formatting
    total = len(prs.slides)
    for index, slide in enumerate(prs.slides):
        if index == 0:
            continue  # no page number on the cover
        _add_textbox(
            slide, name="page_number", text=f"{index}  /  {total - 1}",
            left=_PAGE_NUMBER_X, top=_FOOTER_Y,
            width=_PAGE_NUMBER_WIDTH, height=_FOOTER_HEIGHT,
            font_pt=_FOOTER_PT, colour=_BRAND_LIGHT,
            align=PP_ALIGN.RIGHT,
        )


# ---------------------------------------------------------------------------
# Visual identity passes — typography + accent geometry
# ---------------------------------------------------------------------------


def _apply_typography(prs: Presentation, language: str) -> None:
    """Set Latin + East-Asian font on every run across every slide.

    Default Calibri is the biggest "AI-generated deck" tell. We walk
    every shape post-build and write both ``<a:latin>`` and ``<a:ea>``
    typeface XML on every run — leaving the east-asian slot at the
    PowerPoint default would make CJK chars render in a font that
    doesn't match the Latin choice.
    """
    latin, east_asian = _FONT_FAMILIES.get(language, _DEFAULT_FONT_FAMILY)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = latin
                    if east_asian:
                        _set_east_asian_typeface(run, east_asian)


def _set_east_asian_typeface(run, family: str) -> None:
    """Write the ``<a:ea typeface=...>`` element on a run's rPr.

    python-pptx's ``run.font.name`` setter only writes the Latin
    typeface (``<a:latin>``). PowerPoint consults a SEPARATE
    east-asian slot when laying out CJK code points; this helper
    fills it.
    """
    r_pr = run._r.get_or_add_rPr()
    existing = r_pr.find(qn("a:ea"))
    if existing is not None:
        r_pr.remove(existing)
    ea = r_pr.makeelement(qn("a:ea"), {"typeface": family}, nsmap=None)
    r_pr.append(ea)


def _decorate_with_accents(prs: Presentation) -> None:
    """Place the accent shapes (cover left band, top bar on content slides).

    Idempotent: if the shapes already exist from a previous build pass
    (rare but possible in tests that re-run ``_build``), they're left in
    place — a name match suppresses re-add. The shapes are sent to the
    back of the slide's z-order so they sit BEHIND any text on the slide.
    """
    for index, slide in enumerate(prs.slides):
        if index == 0:
            _add_cover_left_band(slide)
        else:
            _add_top_accent_bar(slide)


def _add_cover_left_band(slide) -> None:
    if _has_named_shape(slide, "accent_left"):
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        _ACCENT_LEFT_WIDTH, _SLIDE_HEIGHT,
    )
    shape.name = "accent_left"
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _BRAND_DARK
    _send_shape_to_back(shape, slide)


def _add_top_accent_bar(slide) -> None:
    if _has_named_shape(slide, "accent_top"):
        return
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Emu(0), Emu(0),
        _SLIDE_WIDTH, _ACCENT_TOP_HEIGHT,
    )
    shape.name = "accent_top"
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = _BRAND_DARK
    _send_shape_to_back(shape, slide)


def _has_named_shape(slide, name: str) -> bool:
    return any(shape.name == name for shape in slide.shapes)


def _send_shape_to_back(shape, slide) -> None:
    """Move ``shape`` to be the first child of ``spTree`` (= back of z-order)."""
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    # spTree's first two children are nvGrpSpPr + grpSpPr (group metadata);
    # everything after that is a shape in z-order. Insert at index 2 so the
    # band lands BEHIND every text shape but the metadata stays intact.
    sp_tree.insert(2, sp)


# ---------------------------------------------------------------------------
# Dark-mode pass (opt-in via ExportOptions.dark_mode)
# ---------------------------------------------------------------------------


def _apply_dark_mode(prs: Presentation) -> None:
    """Swap the light palette for the dark palette on every slide.

    The exporter builds the deck with the light palette unconditionally,
    then this post-pass re-colours individual shapes / runs / table cells
    by looking up their current RGB in the light→dark mapping. The
    approach is intentionally non-invasive — we don't refactor the 100+
    direct ``_BRAND_*`` references into a palette-aware lookup; instead
    we walk the rendered tree after the fact.

    Steps per slide:
    1. Solid-fill the slide background with ``_DARK_SLIDE_BG``.
    2. Walk every shape:
       - If it's a table, iterate the table's cells (fills + text + borders).
       - Otherwise recolour the shape's own fill + text frame.
    """
    for slide in prs.slides:
        _set_slide_background(slide, _DARK_SLIDE_BG)
        for shape in slide.shapes:
            _recolor_shape(shape)


def _set_slide_background(slide, colour: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _recolor_shape(shape) -> None:
    """Single shape: swap its fill, text-run colours, and (if table) its
    per-cell fills + borders + cell-level runs."""
    if shape.has_table:
        for cell in _iter_table_cells(shape.table):
            _swap_fill(cell)
            _swap_text_colors(cell)
            _swap_cell_border_colors(cell)
        return
    _swap_fill(shape)
    if shape.has_text_frame:
        _swap_text_colors(shape)


def _iter_table_cells(table):
    """python-pptx exposes ``iter_cells`` on Table but the API name has
    changed between versions; this wrapper falls through to the manual
    row/col iteration when needed."""
    iter_cells = getattr(table, "iter_cells", None)
    if iter_cells is not None:
        yield from iter_cells()
        return
    for row in table.rows:
        yield from row.cells


def _swap_fill(shape_or_cell) -> None:
    fill = getattr(shape_or_cell, "fill", None)
    if fill is None:
        return
    try:
        rgb = fill.fore_color.rgb
    except (AttributeError, ValueError, TypeError):
        return
    if rgb is None:
        return
    key = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
    new = _LIGHT_TO_DARK_FILL.get(key)
    if new is None:
        return
    fill.solid()
    fill.fore_color.rgb = RGBColor(*new)


def _swap_text_colors(shape_or_cell) -> None:
    """Swap every run's text colour for the dark-mode equivalent.

    Safety net for runs that the builders forgot to colour explicitly:
    when ``font.color.rgb`` is ``None`` (theme inheritance, renders as
    near-black on screen) or pure black ``(0,0,0)``, promote to the
    dark-mode body colour ``#E5E7EB``. Without this fallback such runs
    would render as black-on-dark — invisible. See
    ``.claude/agents/rules/deck-design.md`` "Dark-mode contract".
    """
    text_frame = getattr(shape_or_cell, "text_frame", None)
    if text_frame is None:
        return
    near_white = RGBColor(0xE5, 0xE7, 0xEB)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                rgb = run.font.color.rgb
            except (AttributeError, ValueError, TypeError):
                rgb = None
            if rgb is None or (int(rgb[0]), int(rgb[1]), int(rgb[2])) == (0, 0, 0):
                run.font.color.rgb = near_white
                continue
            key = (int(rgb[0]), int(rgb[1]), int(rgb[2]))
            new = _LIGHT_TO_DARK_TEXT.get(key)
            if new is not None:
                run.font.color.rgb = RGBColor(*new)


def _swap_cell_border_colors(cell) -> None:
    """Walk the cell's ``<a:lnX>`` border elements and recolour any
    ``<a:srgbClr>`` whose value matches the light-palette divider /
    header-rule colours."""
    tc_pr = cell._tc.find(qn("a:tcPr"))
    if tc_pr is None:
        return
    for edge in ("L", "R", "T", "B"):
        ln = tc_pr.find(qn(f"a:ln{edge}"))
        if ln is None:
            continue
        solid = ln.find(qn("a:solidFill"))
        if solid is None:
            continue
        clr = solid.find(qn("a:srgbClr"))
        if clr is None:
            continue
        val = clr.get("val", "")
        if len(val) != 6:
            continue
        try:
            key = (int(val[0:2], 16), int(val[2:4], 16), int(val[4:6], 16))
        except ValueError:
            continue
        new = _LIGHT_TO_DARK_FILL.get(key)
        if new is None:
            continue
        clr.set("val", f"{new[0]:02X}{new[1]:02X}{new[2]:02X}")
