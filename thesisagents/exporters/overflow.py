"""Slide-overflow inspector for ThesisAgents decks (library home).

This is the importable home of the overflow check. ``scripts/check_overflow.py``
is a thin CLI wrapper around it, and ``thesisagents.exporters.review`` folds it
into the one-stop ``review_deck`` audit. Keeping the logic in the package (not a
script) means the MCP ``pptx_review`` tool, the CLI ``review`` subcommand, the
regression test, and the manual script all share one implementation.

What "overflow" means here (mirrors ``slide-overflow-check.md``):

- 16:9 widescreen: slide is 13.333" x 7.5".
- Body sits between ``BODY_TOP = 1.5"`` and ``FOOTER_GUARD = 7.05"`` (the line
  where page numbers / footer copy live). Nothing may render past 7.05".
- A shape overflows when its *wrapped, rendered* text height exceeds either
  (1) the shape's own height, or (2) ``7.05"`` measured from the slide top.

The wrap estimate reads each run's actual ``font.size`` (the exporter sets it
explicitly per run), classifies each character as full-width (CJK / kana / hangul
/ full-width forms ≈ 1.0 em) or half-width (Latin / digits / punctuation ≈ 0.55
em), accumulates width per line, and wraps when a line exceeds the box's inner
width. Line height is the run's font size × 1.2 (PowerPoint single spacing). It is
a *rough* estimate — deliberately conservative — so it catches gross overflow
without needing a font-metrics library.

Importable: ``check_pptx(path) -> list[Violation]`` and
``check_pptx_from_prs(prs) -> list[Violation]``.
"""
from __future__ import annotations

import sys
import unicodedata
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu

_EMU_PER_INCH = 914400
_FOOTER_GUARD_IN = 7.05
_FOOTER_GUARD_EMU = int(_FOOTER_GUARD_IN * _EMU_PER_INCH)
# python-pptx default textbox inner margins are 0.1" left + 0.1" right.
_DEFAULT_SIDE_MARGIN_IN = 0.1
_FULL_WIDTH_EM = 1.0       # CJK / kana / hangul / full-width forms
_HALF_WIDTH_EM = 0.55      # Latin / digits / ASCII punctuation
_LINE_SPACING = 1.2        # PowerPoint single line spacing ≈ 1.2 × font size
_DEFAULT_FONT_PT = 18      # used only when a run carries no explicit size
_TABLE_FONT_PT = 14        # _TABLE_PT in pptx.py — cell font when a run has none
_CELL_V_MARGIN_IN = 0.1    # exporter sets 0.05" top + 0.05" bottom per cell
# Box-overflow tolerance: ignore a sub-fraction-of-a-line overshoot so rounding
# in the estimate doesn't flag a box that visually fits.
_BOX_TOLERANCE_IN = 0.08

# Chrome / decoration the exporter places intentionally — these are NOT body
# content and never "overflow" in the meaningful sense: the top/left accent bars
# are fixed-geometry rectangles, and the page number + footer live *at* the
# footer line (7.05") by design, so a footer-guard check on them is a false
# positive. Everything else (title / body / subhead / kpi / rq_box /
# paper_subtitle / tables / figures) is real content and gets checked.
_CHROME_NAMES = frozenset({"page_number", "footer"})
_CHROME_PREFIXES = ("accent",)


def _is_chrome(name: str) -> bool:
    return name in _CHROME_NAMES or name.startswith(_CHROME_PREFIXES)


@dataclass(frozen=True)
class Violation:
    slide: int
    shape: str
    kind: str            # "overflows its box" | "crosses footer guard"
    rendered_in: float   # measured value, inches
    limit_in: float      # the limit it broke, inches


def _is_full_width(ch: str) -> bool:
    """True for characters that occupy ~1 em (CJK, kana, hangul, full-width)."""
    if ch in ("\t", "\n"):
        return False
    return unicodedata.east_asian_width(ch) in ("F", "W")


def _char_em(ch: str) -> float:
    return _FULL_WIDTH_EM if _is_full_width(ch) else _HALF_WIDTH_EM


def _run_font_pt(run, fallback: int) -> int:
    size = run.font.size
    return int(size.pt) if size is not None else fallback


def _paragraph_lines(paragraph, inner_width_pt: float, fallback_pt: int) -> tuple[int, int]:
    """Estimate (wrapped line count, max font pt) for one paragraph.

    Width is accumulated per character at the run's own font size, so a mixed
    CJK + Latin line wraps where it actually would. An empty paragraph still
    occupies one line at the fallback size.
    """
    runs = list(paragraph.runs)
    if not runs:
        return 1, fallback_pt
    max_pt = fallback_pt
    line_w = 0.0
    lines = 1
    for run in runs:
        pt = _run_font_pt(run, fallback_pt)
        max_pt = max(max_pt, pt)
        for ch in run.text:
            if ch == "\n":
                lines += 1
                line_w = 0.0
                continue
            char_w = _char_em(ch) * pt
            if line_w + char_w > inner_width_pt and line_w > 0:
                lines += 1
                line_w = char_w
            else:
                line_w += char_w
    return lines, max_pt


def _text_height_in(text_frame, box_width_emu: int) -> float:
    """Estimated rendered height of a text frame, in inches.

    An empty text frame (decorative rectangle, blank placeholder) contributes
    no height — it must not be charged a fallback line.
    """
    if not (text_frame.text or "").strip():
        return 0.0
    box_width_in = box_width_emu / _EMU_PER_INCH
    ml = (text_frame.margin_left or Emu(int(_DEFAULT_SIDE_MARGIN_IN * _EMU_PER_INCH)))
    mr = (text_frame.margin_right or Emu(int(_DEFAULT_SIDE_MARGIN_IN * _EMU_PER_INCH)))
    inner_width_in = max(0.1, box_width_in - (ml + mr) / _EMU_PER_INCH)
    inner_width_pt = inner_width_in * 72
    total_pt = 0.0
    for paragraph in text_frame.paragraphs:
        lines, max_pt = _paragraph_lines(paragraph, inner_width_pt, _DEFAULT_FONT_PT)
        total_pt += lines * max_pt * _LINE_SPACING
    return total_pt / 72


def _table_height_in(shape) -> float:
    """Estimated *rendered* height of a table, in inches. python-pptx grows a
    row to fit wrapped cell text, but the GraphicFrame's declared ``height``
    does not change — so a many-row or long-cell table renders far taller than
    declared and can cross the footer guard while ``shape.height`` says it fits.
    Sum each row's height from its tallest cell's wrapped line count.
    """
    table = shape.table
    col_w = [c.width or 0 for c in table.columns]
    total = 0.0
    for r in range(len(table.rows)):
        row_lines = 1
        for c in range(len(table.columns)):
            inner_in = max(0.1, col_w[c] / _EMU_PER_INCH - 2 * _CELL_V_MARGIN_IN)
            inner_pt = inner_in * 72
            cell_lines = sum(
                _paragraph_lines(p, inner_pt, _TABLE_FONT_PT)[0]
                for p in table.cell(r, c).text_frame.paragraphs
            )
            row_lines = max(row_lines, cell_lines)
        total += row_lines * _TABLE_FONT_PT * _LINE_SPACING / 72 + _CELL_V_MARGIN_IN
    return total


def _shape_violations(slide_idx: int, shape) -> list[Violation]:
    out: list[Violation] = []
    name = getattr(shape, "name", "?") or "?"
    if _is_chrome(name):
        return out  # exporter-placed accent bars / page number / footer
    top = shape.top or 0
    height = shape.height or 0
    if getattr(shape, "has_table", False):
        # Use the estimated rendered height (≥ declared) for the footer-guard
        # check, since the table grows past its declared box when cells wrap.
        rendered_in = _table_height_in(shape)
        height = max(height, int(rendered_in * _EMU_PER_INCH))
    if shape.has_text_frame:
        tf = shape.text_frame
        shrink = tf.auto_size == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        if not shrink:
            rendered_in = _text_height_in(tf, shape.width or 0)
            if rendered_in - height / _EMU_PER_INCH > _BOX_TOLERANCE_IN:
                out.append(Violation(
                    slide_idx, name, "overflows its box",
                    round(rendered_in, 2), round(height / _EMU_PER_INCH, 2),
                ))
    # Footer-guard check applies to every content shape (incl. tables, pictures,
    # shrink-to-fit titles): the box itself must clear 7.05".
    bottom = top + height
    if bottom > _FOOTER_GUARD_EMU + 1:
        out.append(Violation(
            slide_idx, name, "crosses footer guard",
            round(bottom / _EMU_PER_INCH, 2), _FOOTER_GUARD_IN,
        ))
    return out


def check_pptx_from_prs(prs) -> list[Violation]:
    """Walk every slide / shape of an open Presentation and return overflow
    violations (empty = clean). Split from ``check_pptx`` so tests can build a
    deck in memory without writing a temp file."""
    violations: list[Violation] = []
    for idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            violations.extend(_shape_violations(idx, shape))
    return violations


def check_pptx(path: str | Path) -> list[Violation]:
    """Walk every slide / shape and return overflow violations (empty = clean)."""
    return check_pptx_from_prs(Presentation(str(path)))


def _report(path: str | Path) -> bool:
    prs = Presentation(str(path))
    n_slides = len(prs.slides)
    n_shapes = sum(len(s.shapes) for s in prs.slides)
    violations = check_pptx_from_prs(prs)
    print(f"overflow check — {path}")
    print(f"slides:        {n_slides}")
    print(f"shapes:        {n_shapes}")
    print(f"violations:    {len(violations)}")
    for v in violations:
        print(f"  slide {v.slide}, shape \"{v.shape}\": {v.kind} "
              f"— rendered {v.rendered_in}\" vs {v.limit_in}\"")
    verdict = "PASS" if not violations else "FAIL"
    print(f"verdict:       {verdict}")
    return not violations


def main(argv: list[str]) -> int:
    if not argv:
        print("usage: check_overflow.py <deck.pptx> [more.pptx ...]")
        return 2
    failed = 0
    for i, path in enumerate(argv):
        if i:
            print()
        if not _report(path):
            failed += 1
    return failed


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
