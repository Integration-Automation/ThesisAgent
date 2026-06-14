"""Regression tests for the canonical overflow inspector (scripts/check_overflow.py).

The inspector is what the `slide-overflow-check` subagent runs, so its calibration
matters: it must flag genuine content overflow, ignore exporter-placed chrome
(accent bars / page number / footer), and treat empty decorative frames as
zero-height. These tests pin that behaviour with synthetic decks so they don't
depend on whether any shipped deck happens to be clean.
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

_SCRIPT = Path(__file__).resolve().parents[1] / "scripts" / "check_overflow.py"


def _load_inspector():
    spec = importlib.util.spec_from_file_location("check_overflow", _SCRIPT)
    mod = importlib.util.module_from_spec(spec)
    # Register before exec so the frozen dataclass can resolve its own module
    # (dataclasses looks the class's module up in sys.modules on 3.14+).
    sys.modules[spec.name] = mod
    spec.loader.exec_module(mod)
    return mod


def _textbox(slide, *, name, text, left, top, width, height, font_pt=19):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.name = name
    tf = box.text_frame
    tf.word_wrap = True
    tf.text = text
    for run in tf.paragraphs[0].runs:
        run.font.size = Pt(font_pt)
    return box


def _blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def test_short_text_in_big_box_is_clean():
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    _textbox(slide, name="body", text="一句話", left=0.5, top=2.0, width=12.0, height=2.0)
    assert mod.check_pptx_from_prs(prs) == []


def test_long_text_in_tiny_box_overflows():
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    # ~200 CJK chars at 19pt in a 0.4"-high box must overflow its box.
    _textbox(slide, name="body", text="字" * 200, left=0.5, top=2.0, width=6.0, height=0.4)
    violations = mod.check_pptx_from_prs(prs)
    assert any(v.kind == "overflows its box" and v.shape == "body" for v in violations)


def test_content_shape_past_footer_guard_is_flagged():
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    _textbox(slide, name="body", text="尾端內容", left=0.5, top=7.2, width=6.0, height=0.5)
    violations = mod.check_pptx_from_prs(prs)
    assert any(v.kind == "crosses footer guard" and v.shape == "body" for v in violations)


def test_chrome_shapes_are_exempt():
    # page_number / footer live at the 7.05" line by design; accent bars are
    # fixed-geometry decoration. None should ever be flagged.
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    _textbox(slide, name="page_number", text="1 / 19", left=11.0, top=7.2, width=2.0, height=0.3)
    _textbox(slide, name="footer", text="footer copy", left=0.5, top=7.2, width=6.0, height=0.3)
    accent = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    accent.name = "accent_left"
    assert mod.check_pptx_from_prs(prs) == []


def test_empty_decorative_frame_is_zero_height():
    # A blank accent rectangle inside the body must not be charged a fallback line.
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.0), Inches(0.08))
    box.name = "subhead"  # non-chrome name, but empty text → no overflow
    box.text_frame.text = ""
    assert mod.check_pptx_from_prs(prs) == []


def test_small_table_is_clean():
    # A few short-cell rows near the top render within the footer guard.
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    tbl = slide.shapes.add_table(4, 3, Inches(0.5), Inches(2.0), Inches(12), Inches(2.0))
    for r in range(4):
        for c in range(3):
            tbl.table.cell(r, c).text = f"r{r}c{c}"
    assert mod.check_pptx_from_prs(prs) == []


def test_long_table_overflow_is_estimated():
    # python-pptx keeps the declared 3.0" height, but 10 rows of wrapping cells
    # render far taller and cross the footer guard. The inspector must estimate
    # the grown height, not trust shape.height.
    mod = _load_inspector()
    prs = Presentation()
    slide = _blank_slide(prs)
    tbl = slide.shapes.add_table(10, 3, Inches(0.5), Inches(2.3), Inches(12), Inches(3.0))
    long = "a long table cell value that wraps across several lines on the slide " * 2
    for r in range(10):
        for c in range(3):
            tbl.table.cell(r, c).text = long
    violations = mod.check_pptx_from_prs(prs)
    assert any(v.kind == "crosses footer guard" for v in violations)


def test_exported_rich_deck_has_no_overflow(tmp_path):
    """End-to-end guard: a rich deck with deliberately long contribution /
    method / pain-point text (the content that used to overflow fixed-height
    boxes) must export overflow-free. This wires the inspector into the suite
    so the adaptive stacked-section + quadrant pagination can't silently
    regress. See pptx.py `_add_stacked_section` / `_pain_points_per_slide`.
    """
    from thesisagents.core.models import (
        ExportOptions,
        Paper,
        PaperCollection,
        PaperSummary,
        Query,
    )
    from thesisagents.exporters.pptx import PptxExporter

    mod = _load_inspector()
    long_body = (
        "VAE(變分自編碼器)編碼器把 prompt 切成對抗潛在向量 $z_a$ 與良性潛在向量 "
        "$z_b$,訓練目標最小化互資訊 $I(z_a;z_b|E_p)$,透過 Data Processing "
        "Inequality(DPI,限制條件互資訊在資料處理鏈中不會上升)保證分離,對改述攻擊具備強穩定性。"
    )
    long_bullets = (
        "規則式過濾遇到改述(攻擊者改寫提示文字)或混淆(以異常編碼或拼字繞過)攻擊就破功",
        "對抗訓練要 fine-tune LLM,正常任務的效能會明顯下降而且部署成本高",
        "每出現一類新攻擊就得再補一批新規則,缺乏對對抗與良性訊號的原理性分離",
    )
    summary = PaperSummary(
        language="zh-tw",
        pain_points=tuple(
            (f"痛點 {i}:既有防禦的根本侷限與成本問題", long_bullets) for i in range(1, 5)
        ),
        research_question="在維持即時延遲且不傷害正常查詢效能的前提下能否事前中和對抗成分?",
        contributions_detailed=tuple(
            (f"{i}. 互資訊式語意分解與譜圖意圖分類", long_body) for i in range(1, 5)
        ),
        core_observation=long_body,
        model="test",
    )
    paper = Paper(
        source="local", source_id="t", title="A Long-Content Thesis Deck for Overflow",
        authors=("Test Author",), year=2026, venue="Test University · 碩士學位論文",
        abstract="", url="", summary=summary,
    )
    collection = PaperCollection(
        query=Query(keywords="overflow", sources=("local",)), papers=(paper,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="overflow-e2e",
        language="zh-tw", dark_mode=True,
    )
    out_path = PptxExporter().export(collection, options)
    violations = mod.check_pptx(out_path)
    assert violations == [], "\n".join(
        f"slide {v.slide} {v.shape}: {v.kind} {v.rendered_in}\" vs {v.limit_in}\""
        for v in violations
    )
