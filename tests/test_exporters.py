"""Render each exporter to tmp_path and re-open the artefact."""

from __future__ import annotations

import json

from autopapertoppt.core.models import ExportOptions, PaperCollection, Query
from autopapertoppt.exporters import export_collection


def _collection(sample_papers) -> PaperCollection:
    query = Query(
        keywords="attention",
        sources=("arxiv",),
        max_results=10,
    )
    return PaperCollection(query=query, papers=tuple(sample_papers))


def test_bibtex_exporter(sample_papers, tmp_path):
    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("bib",),
        out_dir=str(tmp_path),
        filename_stem="test",
    )
    written = export_collection(collection, options)
    bib_text = written["bib"].read_text(encoding="utf-8")
    assert "@" in bib_text
    assert "title = {Sample Paper on Attention}" in bib_text
    assert "\\&" in bib_text
    assert "\\{braces\\}" in bib_text


def test_markdown_exporter(sample_papers, tmp_path):
    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("md",),
        out_dir=str(tmp_path),
        filename_stem="test",
    )
    written = export_collection(collection, options)
    md_text = written["md"].read_text(encoding="utf-8")
    assert md_text.startswith("# Paper search:")
    assert "Sample Paper on Attention" in md_text
    assert "Alice Anderson" in md_text
    assert "**Abstract**" in md_text


def test_markdown_exporter_no_abstract(sample_papers, tmp_path):
    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("md",),
        out_dir=str(tmp_path),
        filename_stem="test",
        include_abstract=False,
    )
    written = export_collection(collection, options)
    md_text = written["md"].read_text(encoding="utf-8")
    assert "**Abstract**" not in md_text


def test_json_exporter_round_trip(sample_papers, tmp_path):
    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("json",),
        out_dir=str(tmp_path),
        filename_stem="test",
    )
    written = export_collection(collection, options)
    data = json.loads(written["json"].read_text(encoding="utf-8"))
    assert data["query"]["keywords"] == "attention"
    assert len(data["papers"]) == 2
    assert data["papers"][0]["title"] == "Sample Paper on Attention"


def _slide_text(slide, name: str) -> str:
    """Return the text of the shape with the given semantic name, or ''.

    The visual-identity pass inserts ``accent_top`` / ``accent_left``
    rectangles BEFORE the text shapes in z-order, so ``shapes[0]`` is
    no longer reliably the title. Tests pin to the project's semantic
    shape names (`title` / `meta` / `body` / `subhead` / `footer` /
    `page_number` / etc.) instead.
    """
    for shape in slide.shapes:
        if shape.name == name and shape.has_text_frame:
            return shape.text_frame.text
    return ""


def test_pptx_exporter_full_deck(sample_papers, tmp_path):
    from pptx import Presentation

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="test",
    )
    written = export_collection(collection, options)
    presentation = Presentation(str(written["pptx"]))
    # Cover + agenda + per-paper(divider + overview + bg + findings) + references.
    # Sample abstracts are short so the optional "Approach" slide is skipped.
    # 1 + 1 + 2 * 4 + 1 = 11.
    assert len(presentation.slides) == 11
    titles = [_slide_text(s, "title") for s in presentation.slides]
    assert any("Paper Review" in t for t in titles)
    assert any(t == "Agenda" for t in titles)
    assert "References" in titles


def test_pptx_exporter_single_paper_skips_agenda_and_divider(sample_papers, tmp_path):
    from pptx import Presentation

    one = sample_papers[:1]
    collection = PaperCollection(
        query=Query(keywords="attention", sources=("arxiv",), max_results=1),
        papers=tuple(one),
    )
    options = ExportOptions(formats=("pptx",), out_dir=str(tmp_path), filename_stem="one")
    written = export_collection(collection, options)
    presentation = Presentation(str(written["pptx"]))
    # cover + overview + bg + findings + references = 5 (short abstract → no approach slide)
    assert len(presentation.slides) == 5
    titles = [_slide_text(s, "title") for s in presentation.slides]
    assert not any(t == "Agenda" for t in titles)
    assert "References" in titles


def _find_run_color(prs, target_rgb: tuple[int, int, int]) -> bool:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb = run.font.color.rgb
                    except (AttributeError, ValueError, TypeError):
                        continue
                    if rgb is not None and tuple(rgb) == target_rgb:
                        return True
    return False


def test_pptx_default_is_dark_mode(sample_papers, tmp_path):
    """``dark_mode`` defaults to True, so an ExportOptions that doesn't
    explicitly pass the field still produces a dark deck.

    Confirms:
    1. Slide background fill is the dark colour (`#12151B`).
    2. At least one run carries the swapped near-white text colour.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="default-dark",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    bg_rgb = list(prs.slides)[0].background.fill.fore_color.rgb
    assert tuple(bg_rgb) == tuple(RGBColor(0x12, 0x15, 0x1B))
    assert _find_run_color(prs, (0xE5, 0xE7, 0xEB)), (
        "no run was re-coloured to the dark-mode near-white text"
    )


def test_pptx_dark_mode_has_no_invisible_runs(sample_papers, tmp_path):
    """Dark-mode regression guard — no text run may end up with
    ``font.color.rgb is None`` or pure black on the dark slide bg.

    A run with no explicit colour inherits the theme's body-text colour
    (renders as near-black) and the dark-mode post-pass cannot map it
    because there's no source RGB to look up. The recolour pass now
    promotes None / black to ``#E5E7EB`` as a safety net; this test
    pins that fallback so a future builder that forgets to set an
    explicit run colour still produces a readable dark deck.
    """
    from pptx import Presentation

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="dark-readability",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    invisible: list[str] = []
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for p_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(paragraph.runs):
                    text = (run.text or "").strip()
                    if not text:
                        continue
                    try:
                        rgb = run.font.color.rgb
                    except (AttributeError, ValueError, TypeError):
                        rgb = None
                    if rgb is None:
                        invisible.append(
                            f"slide {s_idx} shape {shape.name!r} "
                            f"p{p_idx}r{r_idx}: rgb=None text={text[:30]!r}"
                        )
                    elif tuple(rgb) == (0, 0, 0):
                        invisible.append(
                            f"slide {s_idx} shape {shape.name!r} "
                            f"p{p_idx}r{r_idx}: rgb=black text={text[:30]!r}"
                        )
    assert not invisible, (
        "dark-mode deck contains runs with no explicit (or black) "
        "colour — these render invisible on the dark slide bg:\n  "
        + "\n  ".join(invisible[:10])
    )


def _luminance_255(rgb_tuple: tuple[int, int, int]) -> float:
    return 0.2126 * rgb_tuple[0] + 0.7152 * rgb_tuple[1] + 0.0722 * rgb_tuple[2]


_LIGHT_LUMINANCE_THRESHOLD = 0.7 * 255  # > 178


def _shape_fill_rgb(shape) -> tuple[int, int, int] | None:
    try:
        fill_rgb = shape.fill.fore_color.rgb
    except (AttributeError, ValueError, TypeError):
        return None
    if fill_rgb is None:
        return None
    return (int(fill_rgb[0]), int(fill_rgb[1]), int(fill_rgb[2]))


def _scan_shape_for_light_on_light(s_idx, shape, fill_tuple, bad) -> None:
    tf = getattr(shape, "text_frame", None)
    if tf is None:
        return
    for para in tf.paragraphs:
        for run in para.runs:
            if not (run.text or "").strip():
                continue
            try:
                text_rgb = run.font.color.rgb
            except (AttributeError, ValueError, TypeError):
                continue
            if text_rgb is None:
                continue
            text_tuple = (int(text_rgb[0]), int(text_rgb[1]), int(text_rgb[2]))
            if _luminance_255(text_tuple) > _LIGHT_LUMINANCE_THRESHOLD:
                bad.append(
                    f"slide {s_idx} shape {shape.name!r}: "
                    f"fill={fill_tuple} text={text_tuple} text={run.text[:30]!r}"
                )


def test_pptx_dark_mode_no_light_text_on_light_fill(sample_papers, tmp_path):
    """Dark-mode regression guard — failure mode B (light-on-light).

    The previous regression caught text runs with no explicit colour
    (rgb=None, render as black on dark slide bg). This test catches
    the OTHER failure mode: a shape whose fill is light (luminance >
    ~0.7 × 255) but contains text whose colour is ALSO light → text
    disappears INTO the box. The `_RQ_BOX_FILL` (#F3F6FA near-white)
    bug was the cautionary tale: the box stayed near-white in dark
    mode while its text got swapped to near-white via the post-pass.
    """
    from pptx import Presentation

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="dark-contrast",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))

    bad: list[str] = []
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            fill_tuple = _shape_fill_rgb(shape)
            if fill_tuple is None:
                continue
            if _luminance_255(fill_tuple) <= _LIGHT_LUMINANCE_THRESHOLD:
                continue
            _scan_shape_for_light_on_light(s_idx, shape, fill_tuple, bad)
    assert not bad, (
        "dark-mode deck has light-on-light text that disappears into "
        "the shape fill (extend _LIGHT_TO_DARK_FILL to recolour the "
        "fill, OR don't use a near-white fill in light mode):\n  "
        + "\n  ".join(bad[:10])
    )


def test_pptx_no_red_text_runs(sample_papers, tmp_path):
    """The "No red text" contract: ``_BRAND_ACCENT`` (#C0392B) must
    never be written as a run colour. Bold + ``_BRAND_HIGHLIGHT``
    (teal-700 ``#0E7490``) is the approved emphasis pattern for
    headline text (KPI value, RQ question); ``_BRAND_GREY`` is the
    approved pattern for caption / placeholder / chrome text. Red
    font runs read as error / warning in slide-deck conventions and
    pattern-match strongly to AI-generated KPI emphasis ("look at this
    number!"). Banned across light AND dark modes.

    A regression here means a new (or moved) builder added back a
    ``colour=_BRAND_ACCENT`` parameter or wrote
    ``run.font.color.rgb = _BRAND_ACCENT`` directly.
    """
    from pptx import Presentation

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="no-red",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    red = (0xC0, 0x39, 0x2B)
    offenders: list[str] = []
    for s_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            tf = getattr(shape, "text_frame", None)
            if tf is None:
                continue
            for para in tf.paragraphs:
                for run in para.runs:
                    text = (run.text or "").strip()
                    if not text:
                        continue
                    try:
                        rgb = run.font.color.rgb
                    except (AttributeError, ValueError, TypeError):
                        continue
                    if rgb is not None and tuple(rgb) == red:
                        offenders.append(
                            f"slide {s_idx} shape {shape.name!r}: {text[:40]!r}"
                        )
    assert not offenders, (
        "red text (#C0392B) found — use bold + _BRAND_HIGHLIGHT (teal) "
        "for headlines or _BRAND_GREY for captions instead "
        "(deck-design 'No red text' contract):\n  "
        + "\n  ".join(offenders[:10])
    )


def test_pptx_light_mode_keeps_navy_text(sample_papers, tmp_path):
    """``dark_mode=False`` opt-out skips the post-build recolour pass.

    Confirms:
    1. No slide-level background fill is set (or — if set — it isn't
       the dark colour).
    2. At least one run carries the original navy ``_BRAND_DARK``
       (#1F3A66) colour.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="explicit-light",
        dark_mode=False,
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    # No dark slide background.
    try:
        bg_rgb = list(prs.slides)[0].background.fill.fore_color.rgb
    except (AttributeError, ValueError, TypeError):
        bg_rgb = None
    if bg_rgb is not None:
        assert tuple(bg_rgb) != tuple(RGBColor(0x12, 0x15, 0x1B))
    assert _find_run_color(prs, (0x1F, 0x3A, 0x66)), (
        "no run kept the light-palette navy text colour"
    )


def test_pptx_exporter_no_abstract_skips_content_slides(sample_papers, tmp_path):
    from pptx import Presentation

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="no-abs",
        include_abstract=False,
    )
    written = export_collection(collection, options)
    presentation = Presentation(str(written["pptx"]))
    # cover + agenda + per-paper(section divider + overview only) + references
    # = 1+1+2*2+1 = 7
    assert len(presentation.slides) == 7


def test_xlsx_exporter_round_trip(sample_papers, tmp_path):
    from openpyxl import load_workbook

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("xlsx",),
        out_dir=str(tmp_path),
        filename_stem="test",
    )
    written = export_collection(collection, options)
    workbook = load_workbook(str(written["xlsx"]))
    assert "Papers" in workbook.sheetnames
    assert "Query" in workbook.sheetnames
    sheet = workbook["Papers"]
    # Header row + N papers
    assert sheet.max_row == len(sample_papers) + 1
    assert sheet.cell(row=1, column=1).value == "#"
    assert sheet.cell(row=1, column=2).value == "Title"
    assert sheet.cell(row=2, column=2).value == "Sample Paper on Attention"
    assert sheet.cell(row=2, column=3).value == "Alice Anderson, Bob Brown"
    assert sheet.cell(row=2, column=4).value == 2024
    # Column 5 = "Source" (real publication venue), column 6 = "Indexed via"
    # (the metadata fetcher — arxiv, openalex, etc.). This separation matters
    # so users never see "openalex" sitting in the column labelled "Source".
    assert sheet.cell(row=1, column=5).value == "Source"
    assert sheet.cell(row=1, column=6).value == "Indexed via"
    # Second sample paper has venue="NeurIPS 2023" and source="arxiv".
    assert sheet.cell(row=3, column=5).value == "NeurIPS 2023"
    assert sheet.cell(row=3, column=6).value == "arxiv"
    # URL column is column 8 — verify hyperlink applied
    assert sheet.cell(row=2, column=8).hyperlink is not None
    assert sheet.freeze_panes == "A2"


def test_xlsx_exporter_source_column_uses_venue_not_fetcher_name(tmp_path):
    """Regression: xlsx "Source" column shows publication venue, not fetcher.

    Reproduces the user-reported issue: when a paper is fetched via OpenAlex
    but published in IEEE Access, the spreadsheet's "Source" column must say
    "IEEE Access", not "openalex". The fetcher name moves to "Indexed via".
    """
    from openpyxl import load_workbook

    from autopapertoppt.core.models import Paper

    paper = Paper(
        source="openalex",
        source_id="W4400000000",
        title="A Paper Indexed Through OpenAlex But Published in IEEE Access",
        authors=("Demo Author",),
        year=2025,
        venue="IEEE Access",
        abstract="Demo abstract.",
        url="https://doi.org/10.1109/ACCESS.2025.0000000",
        doi="10.1109/ACCESS.2025.0000000",
    )
    collection = _collection([paper])
    options = ExportOptions(
        formats=("xlsx",),
        out_dir=str(tmp_path),
        filename_stem="venue-source",
    )
    written = export_collection(collection, options)
    workbook = load_workbook(str(written["xlsx"]))
    sheet = workbook["Papers"]
    assert sheet.cell(row=2, column=5).value == "IEEE Access"
    assert sheet.cell(row=2, column=6).value == "openalex"


def _assert_no_ellipsis_anywhere(prs) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name in {"footer", "page_number"}:
                continue
            for paragraph in shape.text_frame.paragraphs:
                assert not paragraph.text.endswith("…"), (
                    f"shape {shape.name!r} has truncated text: "
                    f"{paragraph.text!r}"
                )


def _assert_paper_subtitle_uses_venue(prs, expected_venue: str, forbidden_source: str) -> None:
    """The paper_subtitle line must surface the publication venue, not the
    metadata-fetcher source. Papers crawled via OpenAlex / Semantic Scholar
    but published in an IEEE / ACM venue should display the publisher, not
    the aggregator."""
    matches = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "paper_subtitle":
                text = shape.text_frame.text
                if expected_venue in text:
                    matches += 1
                assert forbidden_source not in text, (
                    f"paper_subtitle leaks fetcher source {forbidden_source!r}: "
                    f"{text!r}"
                )
    assert matches > 0, f"expected venue {expected_venue!r} not seen in any subtitle"


def _assert_title_shapes_shrink_to_fit(prs) -> None:
    from pptx.enum.text import MSO_AUTO_SIZE

    seen = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            if shape.name in {"title", "paper_subtitle"}:
                seen += 1
                assert (
                    shape.text_frame.auto_size
                    == MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                ), f"shape {shape.name!r} is missing shrink-to-fit"
    assert seen > 0


def test_pptx_thesis_style_when_rich_summary_attached(sample_papers, tmp_path):
    """When PaperSummary carries rich fields, the deck expands into thesis-style
    slides (pain-points, technique table, RQ result tables, Q&A …)."""
    from pptx import Presentation

    from autopapertoppt.core.models import Paper, PaperSummary, RqResult

    base = sample_papers[0]
    enriched = Paper(
        source=base.source, source_id=base.source_id, title=base.title,
        authors=base.authors, year=base.year, venue=base.venue,
        abstract=base.abstract, url=base.url, doi=base.doi,
        arxiv_id=base.arxiv_id, citation_count=base.citation_count,
        pdf_url=base.pdf_url,
        summary=PaperSummary(
            language="en",
            pain_points=(
                ("Quadrant A", ("a1", "a2")),
                ("Quadrant B", ("b1",)),
            ),
            research_question="Can X be done with Y?",
            contributions_detailed=(
                ("One", "first contribution"),
                ("Two", "second contribution"),
            ),
            headline_metrics=(("score", "0.86", "0.67"),),
            technique_table=(("X", "Role of X"), ("Y", "Role of Y")),
            literature_table=(
                ("Study", "Feature A", "Feature B"),
                ("Ours", "✓", "✓"),
            ),
            system_flow=("step 1", "step 2", "step 3"),
            method_sections=(("M1", ("d1", "d2")),),
            evaluation_sections=(("E1", ("e1",)),),
            research_questions=(("RQ1", "first question"),),
            rq_results=(
                RqResult(
                    rq_id="RQ1", question="first question?",
                    table=(("metric", "ours"), ("a", "1.0")),
                    analysis=("good",),
                ),
            ),
            core_observation="Core insight here.",
            future_work=("more eval",),
            limitations=("small sample",),
        ),
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("arxiv",), max_results=1),
        papers=(enriched,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="rich",
        language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    # cover + overview + pain + contrib + technique + literature + flow +
    # method + evaluation + rqs + 1 rq result + contrib_summary + lim/future + qa + refs
    assert len(prs.slides) >= 14
    titles = [_slide_text(s, "title") for s in prs.slides]
    assert any("Background & Pain Points" in t for t in titles)
    assert any("Key Technologies" in t for t in titles)
    assert any("RQ1" in t for t in titles)
    assert any("Q&A" in t for t in titles)
    # Combined Limitations & Future Work title must not duplicate "Future Work"
    # (regression: en ``section_limitations`` used to be the combined phrase,
    # producing "Limitations & Future Work & Future Work").
    assert any("Limitations & Future Work" in t for t in titles)
    assert not any("Future Work & Future Work" in t for t in titles)
    # Cover slide must not render the same paper title in both the ``title``
    # and ``subtitle`` shapes — regression: ``_english_subtitle`` previously
    # returned ``paper.title`` verbatim and was rendered even when it matched
    # the title above it.
    cover = prs.slides[0]
    title_texts = [
        s.text_frame.text.strip()
        for s in cover.shapes
        if s.has_text_frame and s.name == "title"
    ]
    subtitle_texts = [
        s.text_frame.text.strip()
        for s in cover.shapes
        if s.has_text_frame and s.name == "subtitle"
    ]
    if title_texts and subtitle_texts:
        assert title_texts[0] != subtitle_texts[0]
    # No textbox may end with "…" — word-wrap handles long content; runtime
    # truncation would silently hide information.
    _assert_no_ellipsis_anywhere(prs)
    # Title-area shapes must use MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE so a long
    # title scales down inside the fixed-height box instead of spilling
    # past the horizontal rule beneath.
    _assert_title_shapes_shrink_to_fit(prs)
    # Widescreen 16:9.
    assert abs(prs.slide_width / prs.slide_height - 16 / 9) < 0.01


def _minimal_rich_summary():
    """A PaperSummary with just enough rich content to trigger the
    thesis-style tier (so paper_subtitle shapes get rendered)."""
    from autopapertoppt.core.models import PaperSummary

    return PaperSummary(
        language="en",
        research_question="Trivial test question?",
        contributions_detailed=(("One", "first contribution"),),
    )


def test_pptx_subtitle_prefers_venue_over_fetcher_source(tmp_path):
    """A paper crawled via OpenAlex but published in IEEE must surface the
    IEEE venue on every content slide — never the literal string
    ``openalex``."""
    from pptx import Presentation

    from autopapertoppt.core.models import Paper, PaperCollection, Query

    paper = Paper(
        source="openalex",
        source_id="W12345",
        title="Securing LLM Workloads at Edge",
        authors=("Hassan Karim",),
        year=2025,
        venue="IEEE Access",
        abstract="Edge LLM workload security.",
        url="https://doi.org/10.1109/ACCESS.2025.999",
        doi="10.1109/ACCESS.2025.999",
        pdf_url="https://example.com/x.pdf",
        summary=_minimal_rich_summary(),
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("openalex",), max_results=1),
        papers=(paper,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="venue",
        language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    _assert_paper_subtitle_uses_venue(
        prs, expected_venue="IEEE Access", forbidden_source="openalex"
    )


def test_pptx_rq_slide_shows_verbatim_research_question(tmp_path):
    """The RQ-result slide must surface the verbatim question from
    ``research_questions`` (keyed by rq_id), not just the short label
    that an author may have stored on the RqResult."""
    from pptx import Presentation

    from autopapertoppt.core.models import (
        Paper,
        PaperCollection,
        PaperSummary,
        Query,
        RqResult,
    )

    verbatim = "How effective are LCCT-specific jailbreak attacks?"
    short_label = "LCCT jailbreak ASR vs general LLMs"
    paper = Paper(
        source="arxiv", source_id="x",
        title="LLM Security Paper",
        authors=("Alice",), year=2025, venue="AAAI 2025",
        abstract="x", url="https://example.com/x",
        pdf_url="https://example.com/x.pdf",
        summary=PaperSummary(
            language="en",
            research_questions=(("RQ1", verbatim),),
            rq_results=(
                RqResult(
                    rq_id="RQ1", question=short_label,
                    table=(("metric", "value"), ("acc", "0.99")),
                    analysis=("good",),
                ),
            ),
        ),
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("arxiv",), max_results=1),
        papers=(paper,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rq-verbatim", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    rq_question_texts = [
        shape.text_frame.text
        for slide in prs.slides
        for shape in slide.shapes
        if shape.has_text_frame and shape.name == "rq_question"
    ]
    assert verbatim in rq_question_texts, (
        f"RQ slide must show the verbatim question; got {rq_question_texts!r}"
    )
    # The short label must NOT leak through when a verbatim entry exists.
    assert short_label not in rq_question_texts


def _section_titles(prs) -> list[str]:
    """Pull the 'title' shape text from every slide (in order)."""
    out: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "title":
                out.append(shape.text_frame.text)
                break
    return out


def _rich_paper_factory(**summary_overrides):
    """Build a single-paper PaperCollection with rich summary overrides."""
    from autopapertoppt.core.models import (
        Paper,
        PaperCollection,
        PaperSummary,
        Query,
    )

    summary_kwargs = {
        "language": "en",
        "research_question": "RQ?",
    }
    summary_kwargs.update(summary_overrides)
    paper = Paper(
        source="arxiv", source_id="x", title="T",
        authors=("A",), year=2025, venue="V",
        abstract="x", url="https://x", pdf_url="https://x.pdf",
        summary=PaperSummary(**summary_kwargs),
    )
    return PaperCollection(
        query=Query(keywords="x", sources=("arxiv",), max_results=1),
        papers=(paper,),
    )


def test_pptx_stacks_slide_paginates_beyond_cap(tmp_path):
    """8 contribution stacks must produce 2 contribution slides
    (1/2 and 2/2), not silently drop entries 6-8."""
    from pptx import Presentation

    stacks = tuple((f"Stack {i}", f"Body {i}") for i in range(1, 9))
    collection = _rich_paper_factory(contributions_detailed=stacks)
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="paginate-stacks", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    titles = _section_titles(prs)
    contrib_titles = [t for t in titles if "Contributions" in t]
    # 5 per slide × 2 = 10 capacity for 8 stacks → exactly 2 slides
    assert any("Contributions (1/2)" in t for t in contrib_titles)
    assert any("Contributions (2/2)" in t for t in contrib_titles)


def test_pptx_kpi_slide_paginates_beyond_six(tmp_path):
    """9 KPI metrics must produce 2 Headline Metrics slides."""
    from pptx import Presentation

    metrics = tuple((f"Metric {i}", str(i), "") for i in range(1, 10))
    collection = _rich_paper_factory(headline_metrics=metrics)
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="paginate-kpi", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    titles = _section_titles(prs)
    kpi_titles = [t for t in titles if "Headline Metrics" in t]
    assert any("(1/2)" in t for t in kpi_titles)
    assert any("(2/2)" in t for t in kpi_titles)


def test_pptx_pain_points_paginates_beyond_four(tmp_path):
    """6 pain-point sections must split into 2 slides (4 + 2)."""
    from pptx import Presentation

    pain_points = tuple(
        (f"Quadrant {i}", (f"b{i}.1", f"b{i}.2")) for i in range(1, 7)
    )
    collection = _rich_paper_factory(pain_points=pain_points)
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="paginate-pain", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    titles = _section_titles(prs)
    pp_titles = [t for t in titles if "Background & Pain Points" in t]
    assert any("(1/2)" in t for t in pp_titles)
    assert any("(2/2)" in t for t in pp_titles)


def test_pptx_rq_analysis_keeps_up_to_six_bullets(tmp_path):
    """RQ analysis bullets — six should render, not 3."""
    from pptx import Presentation

    from autopapertoppt.core.models import RqResult

    rq = RqResult(
        rq_id="RQ1",
        question="Real question text",
        table=(("k", "v"), ("a", "1")),
        analysis=tuple(f"analysis bullet {i}" for i in range(1, 7)),
    )
    collection = _rich_paper_factory(
        research_questions=(("RQ1", "Real question text"),),
        rq_results=(rq,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rq-analysis-cap", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    # Find the RQ1 results slide, grab its body bullets.
    seen_bullets: list[str] = []
    for slide in prs.slides:
        title = next(
            (s.text_frame.text for s in slide.shapes
             if s.has_text_frame and s.name == "title"),
            "",
        )
        if "RQ1" not in title:
            continue
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "body":
                seen_bullets.extend(
                    p.text for p in shape.text_frame.paragraphs if p.text.strip()
                )
    assert sum(1 for b in seen_bullets if "analysis bullet" in b) == 6, (
        f"expected 6 analysis bullets to render; got: {seen_bullets}"
    )


def test_pptx_figure_slide_embeds_image_caption_description(tmp_path):
    """A summary with one figure must produce a figure slide carrying the
    image (as a picture shape), the caption in a subhead, and the
    description bullets in the body."""
    from PIL import Image  # type: ignore  # python-pptx requires Pillow already
    from pptx import Presentation

    # Build a 600x300 test image so it's clearly figure-sized.
    img_path = tmp_path / "test-figure.png"
    Image.new("RGB", (600, 300), color=(120, 180, 220)).save(img_path)

    collection = _rich_paper_factory(
        figures=(
            (
                "Architecture overview (Figure 1)",
                str(img_path),
                ("Component A on the left", "Component B on the right"),
            ),
        ),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="figure-slide", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    fig_slides = [
        s for s in prs.slides
        if any(
            sh.has_text_frame and sh.name == "title"
            and "Figure" in sh.text_frame.text
            for sh in s.shapes
        )
    ]
    assert fig_slides, "no Figure slide was rendered"
    fig_slide = fig_slides[0]
    # Picture present
    pic_shapes = [
        sh for sh in fig_slide.shapes
        if getattr(sh, "name", "") == "figure"
    ]
    assert len(pic_shapes) == 1
    # Caption mentions the supplied text
    captions = [
        sh.text_frame.text
        for sh in fig_slide.shapes
        if sh.has_text_frame and sh.name == "subhead"
    ]
    assert any("Architecture overview" in c for c in captions), captions
    # Description bullets rendered
    bullets = [
        p.text for sh in fig_slide.shapes
        if sh.has_text_frame and sh.name == "body"
        for p in sh.text_frame.paragraphs
    ]
    assert any("Component A on the left" in b for b in bullets)
    assert any("Component B on the right" in b for b in bullets)


def test_pptx_paper_table_slide_renders_caption_table_analysis(tmp_path):
    """paper_tables entry produces a Table slide with caption, table, analysis."""
    from pptx import Presentation

    collection = _rich_paper_factory(
        paper_tables=(
            (
                "Vulnerabilities by prompting strategy",
                (
                    ("Strategy", "Total", "Critical"),
                    ("Efficiency", "124", "37"),
                    ("Feature", "158", "29"),
                ),
                ("Feature-focused yields most total vulns", "Efficiency has the most critical"),
            ),
        ),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="paper-table-slide", language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    table_slides = [
        s for s in prs.slides
        if any(
            sh.has_text_frame and sh.name == "title"
            and "Table" in sh.text_frame.text
            for sh in s.shapes
        )
    ]
    assert table_slides
    captions = [
        sh.text_frame.text
        for sh in table_slides[0].shapes
        if sh.has_text_frame and sh.name == "subhead"
    ]
    assert any("Vulnerabilities" in c for c in captions)


def _full_rich_summary_kwargs():
    from autopapertoppt.core.models import RqResult

    return {
        "language": "en",
        "pain_points": (
            ("Q1", ("p1.1", "p1.2")),
            ("Q2", ("p2.1",)),
            ("Q3", ("p3.1",)),
            ("Q4", ("p4.1",)),
        ),
        "research_question": "RQ?",
        "contributions_detailed": (
            ("C1", "first"),
            ("C2", "second"),
            ("C3", "third"),
            ("C4", "fourth"),
        ),
        "headline_metrics": (
            ("m1", "1", ""),
            ("m2", "2", ""),
            ("m3", "3", ""),
            ("m4", "4", ""),
            ("m5", "5", ""),
            ("m6", "6", ""),
        ),
        "technique_table": (("X", "role"),),
        "literature_table": (("Study", "F"), ("Ours", "1")),
        "system_flow": ("a", "b"),
        "method_sections": (
            ("M1", ("d1",)),
            ("M2", ("d2",)),
        ),
        "evaluation_sections": (
            ("E1", ("e1",)),
            ("E2", ("e2",)),
        ),
        "research_questions": (("RQ1", "q1"), ("RQ2", "q2")),
        "rq_results": (
            RqResult(rq_id="RQ1", question="q1", table=(("k", "v"),), analysis=("a",)),
            RqResult(rq_id="RQ2", question="q2", table=(("k", "v"),), analysis=("a",)),
        ),
        "core_observation": "Core.",
        "future_work": ("more",),
        "limitations": ("small",),
    }


def test_pptx_max_slides_trims_to_budget(tmp_path):
    """--max-slides N caps each paper's deck at N slides."""
    from pptx import Presentation

    collection = _rich_paper_factory(**_full_rich_summary_kwargs())
    full_options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rich-full", language="en",
    )
    full = export_collection(collection, full_options)
    full_count = len(Presentation(str(full["pptx"])).slides)
    assert full_count >= 15  # sanity: full rich tier produces lots of slides

    trimmed_options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rich-trim", language="en",
        max_slides_per_paper=10,
    )
    trimmed = export_collection(collection, trimmed_options)
    trimmed_count = len(Presentation(str(trimmed["pptx"])).slides)
    assert trimmed_count <= 10, f"trimmed deck had {trimmed_count} slides"


def test_pptx_max_slides_keeps_essentials(tmp_path):
    """Even with a tight budget, cover / contributions / references must
    survive. Low-priority categories (figure, literature_table, qa) go first."""
    from PIL import Image  # type: ignore
    from pptx import Presentation

    img = tmp_path / "f.png"
    Image.new("RGB", (600, 300), color=(0, 0, 0)).save(img)

    summary_kwargs = _full_rich_summary_kwargs()
    summary_kwargs["figures"] = (
        ("F1", str(img), ("about F1",)),
        ("F2", str(img), ("about F2",)),
    )
    collection = _rich_paper_factory(**summary_kwargs)
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rich-essentials", language="en",
        max_slides_per_paper=8,
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    titles = _section_titles(prs)
    assert len(prs.slides) <= 8
    # Must keep: cover (slide 0, paper title), Contributions, References
    assert any("Contributions" in t for t in titles), titles
    assert any("References" in t for t in titles), titles
    # Figures and Literature Positioning should be dropped first
    assert not any(t.startswith("Figure") for t in titles), titles
    assert not any("Literature Positioning" in t for t in titles), titles


def test_pptx_max_slides_default_is_twentyfive(tmp_path):
    """Without an explicit cap, the deck is held at ≤25 slides — the
    default in ExportOptions. This protects users from accidentally
    shipping a 40-slide deck for a paper with too many figures."""
    from pptx import Presentation

    summary_kwargs = _full_rich_summary_kwargs()
    summary_kwargs["pain_points"] = tuple(
        (f"Q{i}", (f"b{i}.1", f"b{i}.2", f"b{i}.3")) for i in range(1, 9)
    )
    summary_kwargs["contributions_detailed"] = tuple(
        (f"C{i}", f"body {i}") for i in range(1, 13)
    )
    collection = _rich_paper_factory(**summary_kwargs)
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path),
        filename_stem="rich-default", language="en",
    )
    # No explicit max_slides_per_paper → default 25
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    assert len(prs.slides) <= 25, (
        f"default cap should hold deck at ≤25 slides; got {len(prs.slides)}"
    )


def test_pptx_max_slides_zero_disables_cap(tmp_path):
    """``max_slides_per_paper=0`` opts out of trimming — full deck rendered."""
    from PIL import Image  # type: ignore
    from pptx import Presentation

    img = tmp_path / "f.png"
    Image.new("RGB", (600, 300), color=(0, 0, 0)).save(img)

    summary_kwargs = _full_rich_summary_kwargs()
    summary_kwargs["pain_points"] = tuple(
        (f"Q{i}", (f"b{i}.1", f"b{i}.2", f"b{i}.3")) for i in range(1, 9)
    )
    summary_kwargs["contributions_detailed"] = tuple(
        (f"C{i}", f"body {i}") for i in range(1, 13)
    )
    summary_kwargs["figures"] = tuple(
        (f"Figure {i}", str(img), (f"fig{i} bullet",)) for i in range(1, 6)
    )
    summary_kwargs["paper_tables"] = tuple(
        (f"Table {i}", (("k", "v"), ("a", "1")), (f"row{i} note",))
        for i in range(1, 4)
    )
    # Capped run must trim; uncapped run must not.
    collection = _rich_paper_factory(**summary_kwargs)
    capped = export_collection(
        collection,
        ExportOptions(
            formats=("pptx",), out_dir=str(tmp_path),
            filename_stem="rich-capped", language="en",
            max_slides_per_paper=25,
        ),
    )
    uncapped = export_collection(
        collection,
        ExportOptions(
            formats=("pptx",), out_dir=str(tmp_path),
            filename_stem="rich-uncapped", language="en",
            max_slides_per_paper=0,
        ),
    )
    capped_count = len(Presentation(str(capped["pptx"])).slides)
    uncapped_count = len(Presentation(str(uncapped["pptx"])).slides)
    assert capped_count <= 25
    assert uncapped_count > capped_count, (
        f"max_slides_per_paper=0 should produce more slides than the "
        f"cap-25 run: capped={capped_count}, uncapped={uncapped_count}"
    )


def test_pptx_subtitle_falls_back_to_source_when_venue_missing(tmp_path):
    """When venue is None (e.g. arXiv preprint), the source name is the
    natural fallback because arXiv does double duty as publication channel."""
    from pptx import Presentation

    from autopapertoppt.core.models import Paper, PaperCollection, Query

    paper = Paper(
        source="arxiv",
        source_id="2401.00001",
        title="Some Preprint",
        authors=("Anon",),
        year=2025,
        venue=None,
        abstract="abstract",
        url="https://arxiv.org/abs/2401.00001",
        arxiv_id="2401.00001",
        pdf_url="https://arxiv.org/pdf/2401.00001",
        summary=_minimal_rich_summary(),
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("arxiv",), max_results=1),
        papers=(paper,),
    )
    options = ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="preprint",
        language="en",
    )
    written = export_collection(collection, options)
    prs = Presentation(str(written["pptx"]))
    seen_subtitle = False
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and shape.name == "paper_subtitle":
                text = shape.text_frame.text
                seen_subtitle = True
                assert "arxiv" in text
    assert seen_subtitle


def test_xlsx_exporter_no_abstract(sample_papers, tmp_path):
    from openpyxl import load_workbook

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("xlsx",),
        out_dir=str(tmp_path),
        filename_stem="no-abs",
        include_abstract=False,
    )
    written = export_collection(collection, options)
    workbook = load_workbook(str(written["xlsx"]))
    sheet = workbook["Papers"]
    assert sheet.cell(row=2, column=11).value in (None, "")


def test_export_unknown_format_raises(sample_papers, tmp_path):
    import pytest

    from autopapertoppt.core.exceptions import ExportError

    collection = _collection(sample_papers)
    options = ExportOptions(
        formats=("nope",),
        out_dir=str(tmp_path),
    )
    with pytest.raises(ExportError):
        export_collection(collection, options)
