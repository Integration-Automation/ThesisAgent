"""Regression tests for the dark-text auditor (scripts/_audit_dark_text.py).

The auditor is the manual companion to the exporter's dark-mode regression tests —
it must agree with the exporter (a generated dark deck audits clean) and must catch
the invisibility / red-text failure modes on an arbitrary deck (e.g. a hand-made
one the regression tests never see).
"""
from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

_SCRIPT = Path(__file__).resolve().parents[1] / "scripts" / "_audit_dark_text.py"


def _load_auditor():
    spec = importlib.util.spec_from_file_location("_audit_dark_text", _SCRIPT)
    mod = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = mod  # so the frozen dataclass resolves its module
    spec.loader.exec_module(mod)
    return mod


def test_generated_dark_deck_audits_clean(tmp_path):
    from thesisagents.core.models import (
        ExportOptions,
        Paper,
        PaperCollection,
        PaperSummary,
        Query,
    )
    from thesisagents.exporters.pptx import PptxExporter

    mod = _load_auditor()
    summary = PaperSummary(
        language="zh-tw",
        contributions_detailed=(("1. 方法", "以互資訊分解對抗成分,既準又快。"),),
        headline_metrics=(("準確率", "92.3%", "baseline 65.4"),),
        core_observation="把對抗與良性成分在 latent 空間分離,事前中和。",
        model="test",
    )
    paper = Paper(
        source="local", source_id="t", title="Dark Mode Audit Smoke",
        authors=("A",), year=2026, venue="Test", abstract="", url="", summary=summary,
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("local",)), papers=(paper,),
    )
    out = PptxExporter().export(collection, ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="audit-smoke",
        language="zh-tw", dark_mode=True,
    ))
    hard = [i for i in mod.audit_deck(out) if i.hard]
    assert hard == [], "\n".join(f"slide {i.slide} {i.shape}: {i.kind} {i.detail}" for i in hard)


def _audit_one_run(mod, tmp_path, *, name, text, rgb):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.name = name
    box.text_frame.text = text
    run = box.text_frame.paragraphs[0].runs[0]
    run.font.size = Pt(18)
    run.font.color.rgb = rgb
    path = tmp_path / f"{name}.pptx"
    prs.save(str(path))
    return mod.audit_deck(path)


def test_black_run_is_flagged_invisible(tmp_path):
    mod = _load_auditor()
    issues = _audit_one_run(
        mod, tmp_path, name="body", text="黑字在深色背景上看不見", rgb=RGBColor(0, 0, 0)
    )
    assert any(i.kind == "invisible" and i.hard for i in issues)


def test_red_run_is_flagged(tmp_path):
    mod = _load_auditor()
    issues = _audit_one_run(
        mod, tmp_path, name="kpi", text="92.3%", rgb=RGBColor(0xC0, 0x39, 0x2B)
    )
    assert any(i.kind == "red text" and i.hard for i in issues)
