"""End-to-end pptx editor tests: write a real .pptx with the exporter,
then edit it and reopen.

The deck layout is the full presentable structure (cover + agenda +
per-paper section divider + overview + abstract-section slides + references).
We use ``include_abstract=False`` in the test fixture so the deck stays
predictable in size — 7 slides for the two sample papers:

  0: Cover ("Paper Review: attention")
  1: Agenda
  2: Section divider 1 ("Paper 1 of 2")
  3: Overview paper 1 (title = paper 1's title)
  4: Section divider 2 ("Paper 2 of 2")
  5: Overview paper 2 (title = paper 2's title)
  6: References
"""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from autopapertoppt.core.exceptions import ExportError
from autopapertoppt.core.models import ExportOptions, PaperCollection, Query
from autopapertoppt.exporters import export_collection, pptx_edit


@pytest.fixture()
def deck(tmp_path: Path, sample_papers) -> Path:
    collection = PaperCollection(
        query=Query(keywords="attention", sources=("arxiv",), max_results=10),
        papers=tuple(sample_papers),
    )
    options = ExportOptions(
        formats=("pptx",),
        out_dir=str(tmp_path),
        filename_stem="edit-target",
        include_abstract=False,
    )
    written = export_collection(collection, options)
    return written["pptx"]


def test_inspect_returns_full_deck(deck: Path):
    slides = pptx_edit.inspect(deck)
    assert len(slides) == 7
    titles = [s.title for s in slides]
    assert "Paper Review" in titles[0]
    assert titles[1] == "Agenda"
    assert titles[2] == "Paper 1 of 2"
    assert "Sample Paper on Attention" in titles[3]
    assert titles[6] == "References"
    # Each slide carries the semantic shape names.
    assert any(s.name == "title" for s in slides[3].shapes)


def test_update_overview_slide_title_and_body(deck: Path):
    # slide 3 is paper 1's overview — has title, meta, body, footer
    pptx_edit.update_slide(
        deck, 3, title="Renamed paper 1", body="Custom bullet 1"
    )
    slides = pptx_edit.inspect(deck)
    assert slides[3].title == "Renamed paper 1"
    body_shape = next(s for s in slides[3].shapes if s.name == "body")
    assert body_shape.text == "Custom bullet 1"


def test_update_slide_via_shape_updates(deck: Path):
    pptx_edit.update_slide(deck, 0, shape_updates={0: "Brand new cover title"})
    slides = pptx_edit.inspect(deck)
    assert slides[0].title == "Brand new cover title"


def test_update_slide_out_of_range_raises(deck: Path):
    with pytest.raises(ExportError):
        pptx_edit.update_slide(deck, 99, title="x")


def test_update_unknown_named_shape_raises(deck: Path):
    # cover slide has no `meta` shape — wait, it does (subtitle line). Use an
    # actually-absent name: the references slide has no `meta` shape.
    with pytest.raises(ExportError):
        pptx_edit.update_slide(deck, 6, meta="should fail")


def test_delete_slide(deck: Path):
    before = pptx_edit.inspect(deck)
    assert len(before) == 7
    pptx_edit.delete_slide(deck, 5)  # paper-2 overview
    after = pptx_edit.inspect(deck)
    assert len(after) == 6
    # the survivor at slide 5 is now the References slide
    assert after[5].title == "References"


def test_delete_out_of_range(deck: Path):
    with pytest.raises(ExportError):
        pptx_edit.delete_slide(deck, 99)


def test_reorder_slides(deck: Path):
    before_titles = [s.title for s in pptx_edit.inspect(deck)]
    new_order = [6, 0, 1, 2, 3, 4, 5]
    pptx_edit.reorder_slides(deck, new_order)
    after_titles = [s.title for s in pptx_edit.inspect(deck)]
    assert after_titles[0] == before_titles[6]
    assert after_titles[1] == before_titles[0]


def test_reorder_rejects_non_permutation(deck: Path):
    with pytest.raises(ExportError):
        pptx_edit.reorder_slides(deck, [0, 0, 1, 2, 3, 4, 5])
    with pytest.raises(ExportError):
        pptx_edit.reorder_slides(deck, [0, 1])  # wrong length


def test_add_slide_appends(deck: Path):
    pptx_edit.add_slide(deck, title="Extra Slide", body="Extra body")
    slides = pptx_edit.inspect(deck)
    assert len(slides) == 8
    assert slides[-1].title == "Extra Slide"


def test_add_slide_at_position(deck: Path):
    pptx_edit.add_slide(deck, title="Inserted", body="b", position=1)
    slides = pptx_edit.inspect(deck)
    assert slides[1].title == "Inserted"


def test_inspect_missing_file_raises(tmp_path: Path):
    with pytest.raises(ExportError):
        pptx_edit.inspect(tmp_path / "does_not_exist.pptx")


def test_update_saves_to_out_path(deck: Path, tmp_path: Path):
    target = tmp_path / "copy.pptx"
    pptx_edit.update_slide(deck, 3, title="Copied", out_path=target)
    assert target.exists()
    # original unchanged
    original_slides = pptx_edit.inspect(deck)
    assert "Sample Paper on Attention" in original_slides[3].title
    # but copy reflects change. Find title by semantic shape name —
    # accent rectangles inserted by the visual-identity pass sit at
    # `shapes[0]` now.
    copy_pres = Presentation(str(target))
    slide3 = copy_pres.slides[3]
    title_shapes = [
        sh for sh in slide3.shapes
        if sh.name == "title" and sh.has_text_frame
    ]
    assert title_shapes, "slide 3 has no shape named 'title'"
    assert title_shapes[0].text_frame.text == "Copied"
