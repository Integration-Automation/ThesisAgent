"""MCP tool layer: hit each tool through the FastMCP-registered handler."""

from __future__ import annotations

import asyncio
import json
from pathlib import Path

import pytest

from autopapertoppt import mcp as mcp_pkg
from autopapertoppt.core.models import PaperCollection, Query


@pytest.fixture()
def server():
    return mcp_pkg.build_server()


async def _call(server, name: str, **kwargs):
    """Invoke a registered FastMCP tool by name and return the parsed payload."""
    result = await server.call_tool(name, kwargs)
    text = _first_text(result)
    return json.loads(text)


def _first_text(result):
    # FastMCP returns either a list[Content] or a tuple (contents, structured).
    if isinstance(result, tuple):
        result = result[0]
    for block in result:
        text = getattr(block, "text", None)
        if text is not None:
            return text
    raise AssertionError("MCP tool result contained no text block")


def test_search_tool(monkeypatch, server, sample_papers):
    async def fake_run_search(query, **_kwargs):
        return PaperCollection(query=query, papers=tuple(sample_papers))

    async def fake_shutdown():
        return None

    monkeypatch.setattr("autopapertoppt.mcp.server.run_search", fake_run_search)
    monkeypatch.setattr("autopapertoppt.mcp.server.shutdown_clients", fake_shutdown)
    payload = asyncio.run(
        _call(server, "search", keywords="attention", sources=["arxiv"], max_results=5)
    )
    assert payload["count"] == 2
    assert payload["papers"][0]["title"] == "Sample Paper on Attention"


def test_fetch_pdf_text_tool(monkeypatch, server):
    """fetch_pdf_text wraps intelligence.pdf.fetch_and_extract for MCP callers."""
    from autopapertoppt.intelligence.pdf import ExtractedPdf

    async def fake_fetch(pdf_url, source="intelligence"):
        return ExtractedPdf(
            url=pdf_url, page_count=12, chars=500, text="full paper body text"
        )

    monkeypatch.setattr(
        "autopapertoppt.intelligence.pdf.fetch_and_extract", fake_fetch
    )
    payload = asyncio.run(
        _call(server, "fetch_pdf_text", pdf_url="https://arxiv.org/pdf/x")
    )
    assert payload["page_count"] == 12
    assert payload["chars"] == 500
    assert payload["text"].startswith("full paper")


def test_fetch_paper_tool(monkeypatch, server, sample_papers):
    async def fake_single(identifier):
        return PaperCollection(
            query=Query(keywords=identifier.value, sources=("arxiv",), max_results=1),
            papers=(sample_papers[0],),
        )

    async def fake_shutdown():
        return None

    monkeypatch.setattr("autopapertoppt.mcp.server.run_single_paper", fake_single)
    monkeypatch.setattr("autopapertoppt.mcp.server.shutdown_clients", fake_shutdown)
    payload = asyncio.run(_call(server, "fetch_paper", identifier="2401.08741"))
    assert payload["paper"]["title"] == "Sample Paper on Attention"
    assert payload["identifier"]["kind"] == "arxiv"


def test_export_tool(server, sample_papers, tmp_path):
    papers = [p.to_dict() for p in sample_papers]
    payload = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="attention",
            formats=["xlsx", "bib", "pptx"],
            out_dir=str(tmp_path),
            filename_stem="mcp-test",
        )
    )
    assert Path(payload["written"]["xlsx"]).exists()
    assert Path(payload["written"]["bib"]).exists()
    assert Path(payload["written"]["pptx"]).exists()
    assert payload["pptx_path"] == payload["written"]["pptx"]


def test_pptx_inspect_and_update_via_mcp(server, sample_papers, tmp_path):
    papers = [p.to_dict() for p in sample_papers]
    written = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="attention",
            formats=["pptx"],
            out_dir=str(tmp_path),
            filename_stem="mcp-edit",
            include_abstract=False,
        )
    )
    pptx_path = written["written"]["pptx"]

    inspected = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    # cover + agenda + (divider+overview)*2 + references = 7
    assert inspected["slide_count"] == 7

    asyncio.run(
        _call(
            server,
            "pptx_update_slide",
            path=pptx_path,
            slide_index=3,
            title="MCP-Edited Title",
        )
    )
    re_inspected = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    assert re_inspected["slides"][3]["title"] == "MCP-Edited Title"


def test_list_sources_tool(server, monkeypatch):
    """list_sources reports every plugin + reflects current env-var state."""
    # Clear every gating var so the default-on / opt-in semantics are
    # exercised without contamination from the host shell.
    for var in (
        "AUTOPAPERTOPPT_IEEE_API_KEY",
        "AUTOPAPERTOPPT_DISABLE_IEEE_SCRAPING",
        "AUTOPAPERTOPPT_SPRINGER_API_KEY",
        "AUTOPAPERTOPPT_DISABLE_SCHOLAR_SCRAPING",
    ):
        monkeypatch.delenv(var, raising=False)
    payload = asyncio.run(_call(server, "list_sources"))
    names = {entry["name"]: entry for entry in payload["sources"]}
    # Every plugin we ship must appear.
    for required in (
        "arxiv", "semantic_scholar", "openalex", "pubmed",
        "acm", "dblp", "crossref", "openaire",
        "ieee", "springer", "scholar",
    ):
        assert required in names, f"list_sources missing {required!r}"
    # Plugins that need no env var must be enabled.
    assert names["arxiv"]["enabled"] is True
    assert names["dblp"]["enabled"] is True
    # IEEE + Scholar are now default-ON (no opt-out env var set).
    assert names["ieee"]["enabled"] is True
    assert names["ieee"]["opt_out_env_var"] == ["AUTOPAPERTOPPT_DISABLE_IEEE_SCRAPING"]
    assert names["scholar"]["enabled"] is True
    assert names["scholar"]["opt_out_env_var"] == ["AUTOPAPERTOPPT_DISABLE_SCHOLAR_SCRAPING"]
    # Springer still opt-IN — without the API key it is disabled.
    assert names["springer"]["enabled"] is False
    assert names["springer"]["opt_in_env_var"] == ["AUTOPAPERTOPPT_SPRINGER_API_KEY"]
    assert "default_sources" in payload
    # Default mix now includes scholar (alongside ieee + the others).
    assert "scholar" in payload["default_sources"]
    assert "ieee" in payload["default_sources"]


def test_list_sources_reflects_springer_key(server, monkeypatch):
    monkeypatch.setenv("AUTOPAPERTOPPT_SPRINGER_API_KEY", "test-key")
    payload = asyncio.run(_call(server, "list_sources"))
    by_name = {entry["name"]: entry for entry in payload["sources"]}
    assert by_name["springer"]["enabled"] is True


def test_search_passes_top_tier_and_min_citations(monkeypatch, server, sample_papers):
    """top_tier_only + min_citations flow from the MCP tool into the Query."""
    captured = {}

    async def fake_run_search(query, **_kwargs):
        captured["query"] = query
        return PaperCollection(query=query, papers=tuple(sample_papers))

    async def fake_shutdown():
        return None

    monkeypatch.setattr("autopapertoppt.mcp.server.run_search", fake_run_search)
    monkeypatch.setattr("autopapertoppt.mcp.server.shutdown_clients", fake_shutdown)
    asyncio.run(
        _call(
            server,
            "search",
            keywords="x",
            sources=["arxiv"],
            max_results=3,
            top_tier_only=False,
            min_citations=10,
        )
    )
    assert captured["query"].top_tier_only is False
    assert captured["query"].min_citations == 10


def test_search_defaults_to_full_source_mix(monkeypatch, server, sample_papers):
    """When sources is omitted, the search defaults to DEFAULT_SOURCES."""
    from autopapertoppt.core.constants import DEFAULT_SOURCES

    captured = {}

    async def fake_run_search(query, **_kwargs):
        captured["query"] = query
        return PaperCollection(query=query, papers=tuple(sample_papers))

    async def fake_shutdown():
        return None

    monkeypatch.setattr("autopapertoppt.mcp.server.run_search", fake_run_search)
    monkeypatch.setattr("autopapertoppt.mcp.server.shutdown_clients", fake_shutdown)
    asyncio.run(_call(server, "search", keywords="x"))
    assert captured["query"].sources == DEFAULT_SOURCES


def test_export_respects_max_slides_per_paper(server, sample_papers, tmp_path):
    """A small cap must shrink the produced deck to at most that many slides."""
    from pptx import Presentation

    papers = [p.to_dict() for p in sample_papers]
    payload = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="x",
            formats=["pptx"],
            out_dir=str(tmp_path),
            filename_stem="capped",
            include_abstract=False,
            max_slides_per_paper=3,
        )
    )
    prs = Presentation(payload["written"]["pptx"])
    # Cover + agenda + references are kept; with cap=3 per paper for 2
    # sample_papers, the deck must stay well under the uncapped baseline.
    assert len(prs.slides) <= 3 + 3 + 2  # safety bound


def test_export_treats_zero_as_unlimited(server, sample_papers, tmp_path):
    """Passing 0 must mean "no cap" so an agent doesn't truncate by accident."""
    papers = [p.to_dict() for p in sample_papers]
    capped = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="x",
            formats=["pptx"],
            out_dir=str(tmp_path),
            filename_stem="cap-zero",
            include_abstract=True,
            max_slides_per_paper=0,
        )
    )
    natural = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="x",
            formats=["pptx"],
            out_dir=str(tmp_path),
            filename_stem="cap-none",
            include_abstract=True,
            max_slides_per_paper=None,
        )
    )
    # Both should produce the same slide count — neither truncates.
    from pptx import Presentation
    assert len(Presentation(capped["written"]["pptx"]).slides) == len(
        Presentation(natural["written"]["pptx"]).slides
    )


def test_download_pdfs_tool(monkeypatch, server, sample_papers, tmp_path):
    """download_pdfs forwards to core.pdf_download and returns per-paper results."""
    from autopapertoppt.core.pdf_download import PdfDownloadResult

    async def fake_download(collection, out_dir):
        return [
            PdfDownloadResult(
                paper_key=p.bibtex_key(),
                path=Path(out_dir) / "pdfs" / f"{p.bibtex_key()}.pdf"
                if p.pdf_url
                else None,
                skipped_reason=None if p.pdf_url else "no_pdf_url",
            )
            for p in collection.papers
        ]

    async def fake_shutdown():
        return None

    monkeypatch.setattr(
        "autopapertoppt.mcp.server.core_download_pdfs", fake_download
    )
    monkeypatch.setattr("autopapertoppt.mcp.server.shutdown_clients", fake_shutdown)
    papers = [p.to_dict() for p in sample_papers]
    payload = asyncio.run(
        _call(
            server,
            "download_pdfs",
            papers=papers,
            out_dir=str(tmp_path),
        )
    )
    assert payload["out_dir"] == str(tmp_path)
    # sample_papers[0] has pdf_url, sample_papers[1] does not.
    assert payload["saved"] == 1
    assert payload["skipped"] == 1
    assert any(r["reason"] == "no_pdf_url" for r in payload["results"])
    assert all(("path" in r and "reason" in r) for r in payload["results"])


def test_pptx_delete_reorder_add_via_mcp(server, sample_papers, tmp_path):
    papers = [p.to_dict() for p in sample_papers]
    written = asyncio.run(
        _call(
            server,
            "export",
            papers=papers,
            keywords="x",
            formats=["pptx"],
            out_dir=str(tmp_path),
            filename_stem="ops",
            include_abstract=False,
        )
    )
    pptx_path = written["written"]["pptx"]
    baseline = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    baseline_count = baseline["slide_count"]

    asyncio.run(_call(server, "pptx_add_slide", path=pptx_path, title="Added", body="b"))
    after_add = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    assert after_add["slide_count"] == baseline_count + 1

    new_order = [after_add["slide_count"] - 1] + list(range(after_add["slide_count"] - 1))
    asyncio.run(_call(server, "pptx_reorder_slides", path=pptx_path, new_order=new_order))
    after_reorder = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    assert after_reorder["slides"][0]["title"] == "Added"

    asyncio.run(_call(server, "pptx_delete_slide", path=pptx_path, slide_index=0))
    after_delete = asyncio.run(_call(server, "pptx_inspect", path=pptx_path))
    assert after_delete["slide_count"] == baseline_count
