"""Tests for the one-stop deck reviewer (thesisagents.exporters.review).

review_deck folds three audits into one call — overflow, colour contracts, and
paper_rule section completeness — so these pin that each is wired in and that
completeness only gates a thesis-style deck (a lightweight abstract-only deck
must never be failed for legitimately lacking sections).
"""
from __future__ import annotations

import json

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from thesisagents.core.models import (
    ExportOptions,
    Paper,
    PaperCollection,
    PaperSummary,
    Query,
)
from thesisagents.exporters.pptx import PptxExporter
from thesisagents.exporters.review import review_deck


def _export(tmp_path, summary, *, language="zh-tw", stem="review"):
    paper = Paper(
        source="local", source_id="t", title="審片測試論文",
        authors=("A",), year=2026, venue="Test", abstract="一段研究摘要。",
        url="", summary=summary,
    )
    collection = PaperCollection(
        query=Query(keywords="x", sources=("local",)), papers=(paper,),
    )
    return PptxExporter().export(collection, ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem=stem,
        language=language,
    ))


_FULL_THESIS_SUMMARY = PaperSummary(
    language="zh-tw",
    contributions_detailed=(("貢獻一", "提出新方法。"),),
    technique_table=(("互資訊", "分解對抗成分"),),
    method_sections=(("步驟", ("先分離成分。", "再中和。")),),
    evaluation_sections=(("資料集", ("用三組基準。",)),),
    headline_metrics=(("準確率", "92.3%", "65.4"),),
    core_observation="把對抗與良性成分在 latent 空間分離。",
    model="test",
)


def test_review_clean_thesis_deck_ok(tmp_path):
    out = _export(tmp_path, _FULL_THESIS_SUMMARY)
    review = review_deck(out)
    assert review.language == "zh-tw", "language should auto-detect from titles"
    assert review.thesis_style is True
    assert review.overflow == ()
    assert review.hard_contrast == []
    assert review.missing_sections == (), f"unexpected gaps: {review.missing_sections}"
    assert review.ok is True


def test_review_explicit_language_override(tmp_path):
    out = _export(tmp_path, _FULL_THESIS_SUMMARY)
    review = review_deck(out, language="zh-tw")
    assert review.language == "zh-tw"
    assert review.ok is True


def test_review_flags_overflow(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # A real content shape whose box bottom (top 7.0" + height 1.0") sits past the
    # 7.05" footer guard.
    box = slide.shapes.add_textbox(Inches(1), Inches(7.0), Inches(6), Inches(1.0))
    box.name = "body"
    box.text_frame.text = "這段內容越過了頁尾守線"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x66)
    path = tmp_path / "overflow.pptx"
    prs.save(str(path))

    review = review_deck(path)
    assert any(v.kind == "crosses footer guard" for v in review.overflow)
    assert review.ok is False


def test_review_flags_invisible_run(tmp_path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    box.name = "body"
    box.text_frame.text = "黑字在深色背景上看不見"
    box.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
    box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
    path = tmp_path / "invisible.pptx"
    prs.save(str(path))

    review = review_deck(path)
    assert any(i.kind == "invisible" and i.hard for i in review.contrast)
    assert review.ok is False


def test_review_lightweight_deck_not_failed_for_missing_sections(tmp_path):
    # No summary -> lightweight abstract-only deck. It legitimately lacks most
    # canonical sections, but must NOT be failed for it.
    out = _export(tmp_path, None, stem="lightweight")
    review = review_deck(out)
    assert review.thesis_style is False
    assert review.missing_sections, "a lightweight deck is expected to lack sections"
    assert review.ok is True, "completeness must not gate a lightweight deck"


def test_review_thesis_deck_missing_section_fails(tmp_path):
    # A thesis-style deck (has metrics + core observation) but no methodology or
    # literature slides -> those sections are missing and the deck FAILS.
    partial = PaperSummary(
        language="zh-tw",
        headline_metrics=(("準確率", "92.3%", "65.4"),),
        core_observation="一段核心觀察。",
        model="test",
    )
    out = _export(tmp_path, partial, stem="partial")
    review = review_deck(out)
    assert review.thesis_style is True
    assert "methodology" in review.missing_sections
    assert "literature_review" in review.missing_sections
    assert review.ok is False


def _export_multi(tmp_path, n=2):
    papers = tuple(
        Paper(
            source="arxiv", source_id=f"p{i}", title=f"Paper {i} on Attention",
            authors=("A",), year=2024, venue="V", abstract="Some abstract text.",
            url=f"https://e.com/{i}", arxiv_id=f"24{i:02d}.0000{i}",
        )
        for i in range(1, n + 1)
    )
    collection = PaperCollection(
        query=Query(keywords="attention", sources=("arxiv",)), papers=papers,
    )
    return PptxExporter().export(collection, ExportOptions(
        formats=("pptx",), out_dir=str(tmp_path), filename_stem="multi", language="en",
    ))


def test_review_multipaper_deck_has_references(tmp_path):
    # A real multi-paper deck carries an agenda + references slide; references is
    # therefore present and not flagged.
    out = _export_multi(tmp_path)
    review = review_deck(out)
    assert review.references_missing is False
    assert "references" not in review.missing_sections


def test_review_multipaper_missing_references_fails(tmp_path):
    # A multi-paper-shaped deck (has an Agenda slide) with NO references slide:
    # references is a genuine gap, so the deck FAILS even though it isn't
    # thesis-style.
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for title in ("Cover", "Agenda", "Paper 1 on Attention", "Paper 2 on Attention"):
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(Inches(1), Inches(0.4), Inches(8), Inches(1))
        box.name = "title"
        box.text_frame.text = title
        box.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
        box.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0x1F, 0x3A, 0x66)
    path = tmp_path / "no_refs.pptx"
    prs.save(str(path))

    review = review_deck(path, language="en")
    assert review.references_missing is True
    assert "references" in review.missing_sections
    assert review.ok is False


def test_review_single_paper_rich_deck_does_not_gate_references(tmp_path):
    # A single-paper rich deck (no agenda) folds references into the cover, so a
    # missing references slide must NOT be flagged.
    out = _export(tmp_path, _FULL_THESIS_SUMMARY, stem="single")
    review = review_deck(out)
    assert review.references_missing is False
    assert "references" not in review.missing_sections


def test_review_json_output(tmp_path, capsys):
    from thesisagents.exporters import review as review_mod

    out = _export(tmp_path, _FULL_THESIS_SUMMARY, stem="jsondeck")
    code = review_mod.main(["--json", str(out)])
    payload = json.loads(capsys.readouterr().out)
    assert isinstance(payload, list) and len(payload) == 1
    entry = payload[0]
    assert set(entry) >= {
        "path", "language", "ok", "overflow", "contrast",
        "missing_sections", "references_missing",
    }
    assert entry["ok"] is True
    assert code == 0
