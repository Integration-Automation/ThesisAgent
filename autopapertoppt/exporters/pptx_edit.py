"""Editing operations for an existing .pptx file.

These wrap python-pptx with semantic helpers (`title` / `meta` / `body` lookup
by shape name when the file was written by `PptxExporter`, falling back to
positional access otherwise) and the slide-list XML manipulation required to
delete and reorder slides — python-pptx does not expose those operations
directly.
"""

from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from autopapertoppt.core.exceptions import ExportError

_DEFAULT_TITLE_PT = 28
_DEFAULT_BODY_PT = 14
_SLIDE_LEFT = Inches(0.5)
_SLIDE_TOP = Inches(0.6)
_SLIDE_WIDTH = Inches(9.0)
_TITLE_HEIGHT = Inches(1.2)
_BODY_HEIGHT = Inches(5.4)


@dataclass(frozen=True, slots=True)
class ShapeView:
    """A read-only view of one text-bearing shape on a slide."""

    index: int
    name: str
    text: str


@dataclass(frozen=True, slots=True)
class SlideView:
    """A read-only view of one slide and its text-bearing shapes."""

    index: int
    title: str
    shapes: tuple[ShapeView, ...]


def inspect(path: str | Path) -> list[SlideView]:
    """Return a list of SlideView for `path` so callers can decide what to edit."""
    presentation = _open(path)
    views: list[SlideView] = []
    for slide_index, slide in enumerate(presentation.slides):
        shapes: list[ShapeView] = []
        title = ""
        for shape_index, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text
            shape_name = (shape.name or "").lower()
            shapes.append(ShapeView(index=shape_index, name=shape.name, text=text))
            if not title and (shape_name == "title" or shape_index == 0):
                title = text
        views.append(SlideView(index=slide_index, title=title, shapes=tuple(shapes)))
    return views


def update_slide(
    path: str | Path,
    slide_index: int,
    *,
    title: str | None = None,
    body: str | None = None,
    meta: str | None = None,
    shape_updates: dict[int, str] | None = None,
    out_path: str | Path | None = None,
) -> Path:
    """Update text on `slide_index` and save.

    `title` / `body` / `meta` look the shape up by name when present; if the
    file was produced by PptxExporter that always works. `shape_updates` is the
    fallback — `{shape_index: new_text}` for arbitrary shape edits.
    """
    presentation = _open(path)
    slide = _slide_at(presentation, slide_index)
    by_name = {(s.name or "").lower(): s for s in slide.shapes if s.has_text_frame}
    if title is not None:
        _set_text_by_name(by_name, "title", title, _DEFAULT_TITLE_PT, bold=True)
    if meta is not None:
        _set_text_by_name(by_name, "meta", meta, _DEFAULT_BODY_PT)
    if body is not None:
        _set_text_by_name(by_name, "body", body, _DEFAULT_BODY_PT)
    if shape_updates:
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        for shape_idx, new_text in shape_updates.items():
            if shape_idx < 0 or shape_idx >= len(text_shapes):
                raise ExportError(
                    "pptx_edit",
                    f"shape index {shape_idx} out of range for slide {slide_index}",
                )
            _write_text_frame(text_shapes[shape_idx], new_text)
    return _save(presentation, path, out_path)


def delete_slide(
    path: str | Path,
    slide_index: int,
    out_path: str | Path | None = None,
) -> Path:
    """Delete the slide at `slide_index` and save."""
    presentation = _open(path)
    _ = _slide_at(presentation, slide_index)
    _remove_slide_at(presentation, slide_index)
    return _save(presentation, path, out_path)


def reorder_slides(
    path: str | Path,
    new_order: list[int],
    out_path: str | Path | None = None,
) -> Path:
    """Reorder slides so `new_order[i]` becomes the new index `i`."""
    presentation = _open(path)
    slide_count = len(presentation.slides)
    _validate_permutation(new_order, slide_count)
    sld_id_lst = presentation.slides._sldIdLst  # noqa: SLF001  # python-pptx exposes no public reorder API
    children = list(sld_id_lst)
    reordered = [children[i] for i in new_order]
    for child in children:
        sld_id_lst.remove(child)
    for child in reordered:
        sld_id_lst.append(child)
    return _save(presentation, path, out_path)


def add_slide(
    path: str | Path,
    *,
    title: str,
    body: str = "",
    meta: str = "",
    position: int | None = None,
    out_path: str | Path | None = None,
) -> Path:
    """Append (or insert at `position`) a new slide with title/meta/body textboxes."""
    presentation = _open(path)
    layout = presentation.slide_layouts[6]
    slide = presentation.slides.add_slide(layout)
    _add_named_textbox(
        slide, "title", title, _SLIDE_TOP, _TITLE_HEIGHT, _DEFAULT_TITLE_PT, bold=True
    )
    if meta:
        _add_named_textbox(slide, "meta", meta, Inches(1.8), Inches(0.8), _DEFAULT_BODY_PT)
    if body:
        _add_named_textbox(slide, "body", body, Inches(2.8), _BODY_HEIGHT, _DEFAULT_BODY_PT)
    if position is not None:
        target = max(0, min(position, len(presentation.slides) - 1))
        current_index = len(presentation.slides) - 1
        _move_slide(presentation, current_index, target)
    return _save(presentation, path, out_path)


def _open(path: str | Path) -> Presentation:
    resolved = Path(path).expanduser().resolve()
    if not resolved.exists():
        raise ExportError("pptx_edit", f"pptx file not found: {resolved}")
    return Presentation(str(resolved))


def _save(
    presentation: Presentation,
    original: str | Path,
    out_path: str | Path | None,
) -> Path:
    target = Path(out_path).expanduser().resolve() if out_path else Path(original).resolve()
    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))
    return target


def _slide_at(presentation: Presentation, slide_index: int):
    slides = presentation.slides
    if slide_index < 0 or slide_index >= len(slides):
        raise ExportError(
            "pptx_edit",
            f"slide index {slide_index} out of range (have {len(slides)})",
        )
    return slides[slide_index]


def _set_text_by_name(by_name, key, text, font_pt, bold=False):
    shape = by_name.get(key)
    if shape is None:
        raise ExportError(
            "pptx_edit",
            f"no shape named {key!r} on this slide — use shape_updates instead",
        )
    _write_text_frame(shape, text, font_pt=font_pt, bold=bold)


def _write_text_frame(shape, text: str, *, font_pt: int | None = None, bold: bool = False) -> None:
    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if font_pt is not None:
                run.font.size = Pt(font_pt)
            run.font.bold = bold


def _add_named_textbox(slide, name: str, text: str, top, height, font_pt: int, bold: bool = False):
    box = slide.shapes.add_textbox(_SLIDE_LEFT, top, _SLIDE_WIDTH, height)
    box.name = name
    _write_text_frame(box, text, font_pt=font_pt, bold=bold)
    return box


def _remove_slide_at(presentation: Presentation, slide_index: int) -> None:
    slides = presentation.slides
    sld_id_lst = slides._sldIdLst  # noqa: SLF001
    children = list(sld_id_lst)
    target = children[slide_index]
    rid = target.attrib.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    sld_id_lst.remove(target)
    if rid:
        presentation.part.drop_rel(rid)


def _move_slide(presentation: Presentation, src_index: int, dst_index: int) -> None:
    sld_id_lst = presentation.slides._sldIdLst  # noqa: SLF001
    children = list(sld_id_lst)
    moving = children.pop(src_index)
    children.insert(dst_index, moving)
    for child in list(sld_id_lst):
        sld_id_lst.remove(child)
    for child in children:
        sld_id_lst.append(child)


def _validate_permutation(order: list[int], slide_count: int) -> None:
    if len(order) != slide_count:
        raise ExportError(
            "pptx_edit",
            f"new_order length {len(order)} does not match slide count {slide_count}",
        )
    if sorted(order) != list(range(slide_count)):
        raise ExportError(
            "pptx_edit",
            f"new_order must be a permutation of [0..{slide_count - 1}]",
        )


